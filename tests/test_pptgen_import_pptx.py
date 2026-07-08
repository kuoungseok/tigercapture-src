from __future__ import annotations


def test_import_pptx_deck_reads_text_table_and_image(tmp_path):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    from app.pptgen.import_pptx import import_pptx_deck

    image_path = tmp_path / "hero.png"
    Image.new("RGB", (64, 32), (40, 120, 220)).save(image_path)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    text.text = "Imported Title"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    slide.shapes.add_picture(str(image_path), Inches(6), Inches(1), Inches(2), Inches(1))

    pptx_path = tmp_path / "source.pptx"
    prs.save(pptx_path)

    deck = import_pptx_deck(pptx_path, asset_dir=tmp_path / "assets")
    assert deck.title == "source"
    assert len(deck.slides) == 1
    kinds = {element.kind for element in deck.slides[0].elements}
    assert {"text", "table", "image"} <= kinds
    imported_text = next(element for element in deck.slides[0].elements if element.kind == "text")
    assert imported_text.text == "Imported Title"
    imported_table = next(element for element in deck.slides[0].elements if element.kind == "table")
    assert imported_table.metadata["cells"][1][1] == "2"
    imported_image = next(element for element in deck.slides[0].elements if element.kind == "image")
    assert imported_image.source_path
    assert deck.assets and deck.assets[0]["source"] == "pptx_import"
