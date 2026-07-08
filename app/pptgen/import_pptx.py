"""Best-effort PPTX import for the user PPT generator.

This is intentionally conservative: it imports common PowerPoint content into
TigerCapture's editable deck model without promising full Office fidelity.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from app.pptgen.assets import add_deck_asset
from app.pptgen.schema import DeckSpec, ElementStyle, SlideElement, SlideSpec


SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5


def _require_pptx():
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.enum.text import PP_ALIGN
    except Exception as exc:  # pragma: no cover - dependency guard
        raise RuntimeError("python-pptx is required for PPTX import") from exc
    return Presentation, MSO_SHAPE_TYPE, PP_ALIGN


def _norm(value: Any, total: Any, default: float = 0.0) -> float:
    try:
        denom = max(1.0, float(total))
        number = float(value) / denom
    except Exception:
        number = float(default)
    return max(0.0, min(1.0, number))


def _font_size(paragraph) -> int:
    for run in getattr(paragraph, "runs", []) or []:
        size = getattr(getattr(run, "font", None), "size", None)
        if size is not None:
            try:
                return max(1, int(round(float(size.pt))))
            except Exception:
                pass
    return 28


def _font_family(paragraph, fallback: str) -> str:
    for run in getattr(paragraph, "runs", []) or []:
        name = getattr(getattr(run, "font", None), "name", None)
        if name:
            return str(name)
    return fallback


def _font_flag(paragraph, flag: str) -> bool:
    for run in getattr(paragraph, "runs", []) or []:
        value = getattr(getattr(run, "font", None), flag, None)
        if value is not None:
            return bool(value)
    return False


def _rgb_to_hex(value: Any, fallback: str = "#182033") -> str:
    if value is None:
        return fallback
    try:
        return f"#{int(value[0]):02X}{int(value[1]):02X}{int(value[2]):02X}"
    except Exception:
        pass
    raw = str(value or "").strip()
    if len(raw) == 6:
        return f"#{raw.upper()}"
    return fallback


def _text_color(paragraph, fallback: str = "#182033") -> str:
    for run in getattr(paragraph, "runs", []) or []:
        try:
            rgb = run.font.color.rgb
        except Exception:
            rgb = None
        if rgb is not None:
            return _rgb_to_hex(rgb, fallback)
    return fallback


def _paragraph_align(paragraph, PP_ALIGN) -> str:
    align = getattr(paragraph, "alignment", None)
    if align == PP_ALIGN.CENTER:
        return "center"
    if align == PP_ALIGN.RIGHT:
        return "right"
    return "left"


def _shape_fill(shape, fallback: str = "#F7F9FC") -> str:
    try:
        rgb = shape.fill.fore_color.rgb
    except Exception:
        rgb = None
    return _rgb_to_hex(rgb, fallback)


def _shape_stroke(shape, fallback: str = "#B8C2D6") -> str:
    try:
        rgb = shape.line.color.rgb
    except Exception:
        rgb = None
    return _rgb_to_hex(rgb, fallback)


def _element_rect(shape, slide_width: int, slide_height: int) -> tuple[float, float, float, float]:
    return (
        _norm(getattr(shape, "left", 0), slide_width),
        _norm(getattr(shape, "top", 0), slide_height),
        _norm(getattr(shape, "width", 1), slide_width, 0.2),
        _norm(getattr(shape, "height", 1), slide_height, 0.1),
    )


def _import_table(shape, element_id: str, slide_width: int, slide_height: int) -> SlideElement:
    x, y, w, h = _element_rect(shape, slide_width, slide_height)
    table = shape.table
    rows = len(table.rows)
    cols = len(table.columns)
    cells: list[list[str]] = []
    for r in range(rows):
        row: list[str] = []
        for c in range(cols):
            row.append(str(table.cell(r, c).text or ""))
        cells.append(row)
    return SlideElement(
        id=element_id,
        kind="table",
        name=str(getattr(shape, "name", "") or "Table"),
        x=x,
        y=y,
        w=w,
        h=h,
        style=ElementStyle(fill="#FFFFFF", stroke="#B8C2D6", stroke_width=1.0, color="#182033", font_size=16),
        metadata={"rows": rows, "cols": cols, "header": rows > 1, "cells": cells},
    )


def _import_text(shape, element_id: str, slide_width: int, slide_height: int, fallback_font: str) -> SlideElement:
    x, y, w, h = _element_rect(shape, slide_width, slide_height)
    text = str(getattr(shape, "text", "") or "")
    paragraphs = list(getattr(shape.text_frame, "paragraphs", []) or [])
    paragraph = paragraphs[0] if paragraphs else None
    style = ElementStyle(font_family=fallback_font, font_size=28, color="#182033")
    if paragraph is not None:
        style.font_family = _font_family(paragraph, fallback_font)
        style.font_size = _font_size(paragraph)
        style.bold = _font_flag(paragraph, "bold")
        style.italic = _font_flag(paragraph, "italic")
        style.underline = _font_flag(paragraph, "underline")
        style.color = _text_color(paragraph)
        Presentation, MSO_SHAPE_TYPE, PP_ALIGN = _require_pptx()
        style.align = _paragraph_align(paragraph, PP_ALIGN)
    try:
        if getattr(shape.fill, "type", None) is not None:
            style.fill = _shape_fill(shape, "")
    except Exception:
        pass
    return SlideElement.text_box(
        element_id,
        text,
        x=x,
        y=y,
        w=w,
        h=h,
        font_size=style.font_size,
        font_family=style.font_family,
        bold=style.bold,
        italic=style.italic,
        underline=style.underline,
        color=style.color,
        align=style.align,
    )


def _import_picture(
    deck: DeckSpec,
    shape,
    element_id: str,
    slide_width: int,
    slide_height: int,
    asset_dir: Path,
) -> SlideElement:
    x, y, w, h = _element_rect(shape, slide_width, slide_height)
    image = shape.image
    ext = str(getattr(image, "ext", "png") or "png").lower().lstrip(".")
    asset_dir.mkdir(parents=True, exist_ok=True)
    target = asset_dir / f"{element_id}.{ext}"
    target.write_bytes(image.blob)
    asset = add_deck_asset(deck, target, kind="image", name=str(getattr(shape, "name", "") or target.stem), source="pptx_import")
    element = SlideElement.image(
        element_id,
        target,
        x=x,
        y=y,
        w=w,
        h=h,
        kind="image",
        name=str(getattr(shape, "name", "") or target.stem),
    )
    element.metadata["ppt_asset_id"] = str(asset.get("id") or "")
    element.metadata["source"] = "pptx_import"
    return element


def _import_shape(shape, element_id: str, slide_width: int, slide_height: int) -> SlideElement:
    x, y, w, h = _element_rect(shape, slide_width, slide_height)
    return SlideElement(
        id=element_id,
        kind="shape",
        name=str(getattr(shape, "name", "") or "Shape"),
        x=x,
        y=y,
        w=w,
        h=h,
        style=ElementStyle(fill=_shape_fill(shape), stroke=_shape_stroke(shape), stroke_width=1.0),
        metadata={"source": "pptx_import"},
    )


def import_pptx_deck(path: str | Path, *, asset_dir: str | Path | None = None) -> DeckSpec:
    """Import text, tables, images, and simple shapes from a PPTX file."""
    Presentation, MSO_SHAPE_TYPE, PP_ALIGN = _require_pptx()
    source = Path(path)
    prs = Presentation(str(source))
    deck = DeckSpec(id=f"pptx-{source.stem}", title=source.stem or "Imported PPTX")
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    assets = Path(asset_dir) if asset_dir else source.with_suffix("").parent / f"{source.stem}_pptx_assets"

    for slide_index, ppt_slide in enumerate(prs.slides, start=1):
        slide = SlideSpec(id=f"slide-{slide_index:03d}", title=f"Slide {slide_index}", duration_ms=5000)
        for shape_index, shape in enumerate(ppt_slide.shapes, start=1):
            element_id = f"{slide.id}-el-{shape_index:03d}"
            element: SlideElement | None = None
            if bool(getattr(shape, "has_table", False)):
                element = _import_table(shape, element_id, slide_width, slide_height)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                element = _import_picture(deck, shape, element_id, slide_width, slide_height, assets)
            elif bool(getattr(shape, "has_text_frame", False)) and str(getattr(shape, "text", "") or "").strip():
                element = _import_text(shape, element_id, slide_width, slide_height, deck.theme.font_family)
            elif hasattr(shape, "fill") and hasattr(shape, "line"):
                element = _import_shape(shape, element_id, slide_width, slide_height)
            if element is None:
                continue
            element.z_index = len(slide.elements)
            slide.add_element(element)
        deck.slides.append(slide)

    if not deck.slides:
        deck.slides.append(SlideSpec(id="slide-001", title="Imported PPTX"))
    deck.metadata["source_pptx"] = str(source)
    deck.metadata["importer"] = "python-pptx-basic"
    return deck


__all__ = ["import_pptx_deck"]
