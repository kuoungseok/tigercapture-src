"""PowerPoint-compatible PPTX writer using python-pptx plus timing patches."""
from __future__ import annotations

import shutil
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from app.pptgen.animations import animation_is_active, animation_sequence_sort_key
from app.pptgen.formula import evaluate_numeric_formula, format_formula_value
from app.pptgen.schema import DeckSpec, SlideElement
from app.pptgen.writer_ooxml import _animation_effect_xml, _hex, _timing_xml


SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5
ACTOR_PLACEHOLDER_KINDS = {"video_actor", "ar_pbr_actor", "vrm_actor", "mmd_actor", "audio_actor", "media_actor"}


def _require_pptx():
    try:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception as exc:  # pragma: no cover - dependency guard
        raise RuntimeError("python-pptx is required for PowerPoint-compatible PPTX export") from exc
    return Presentation, ChartData, RGBColor, MSO_CONNECTOR, MSO_SHAPE, PP_ALIGN, XL_CHART_TYPE, XL_LEGEND_POSITION, Inches, Pt


def _left(value: float, Inches):
    return Inches(SLIDE_W_IN * max(0.0, min(1.0, float(value))))


def _top(value: float, Inches):
    return Inches(SLIDE_H_IN * max(0.0, min(1.0, float(value))))


def _width(value: float, Inches):
    return Inches(SLIDE_W_IN * max(0.001, min(1.0, float(value))))


def _height(value: float, Inches):
    return Inches(SLIDE_H_IN * max(0.001, min(1.0, float(value))))


def _rgb(value: str | None, RGBColor, fallback: str = "#FFFFFF"):
    raw = _hex(value, fallback)
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _set_shape_name(shape, element: SlideElement) -> None:
    try:
        shape.name = element.id
    except Exception:
        pass


def _set_fill_and_line(shape, element: SlideElement, deck: DeckSpec, RGBColor) -> None:
    fill_color = element.style.fill or deck.theme.surface
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_color, RGBColor)
    except Exception:
        pass
    try:
        if float(element.style.stroke_width or 0) <= 0:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = _rgb(element.style.stroke or deck.theme.accent, RGBColor)
            shape.line.width = int(max(1, float(element.style.stroke_width or 1)) * 12700)
    except Exception:
        pass


def _paragraph_align(value: str | None, PP_ALIGN):
    normalized = str(value or "").lower()
    if normalized == "center":
        return PP_ALIGN.CENTER
    if normalized == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def _add_text(slide, element: SlideElement, deck: DeckSpec, RGBColor, PP_ALIGN, Inches, Pt):
    shape = slide.shapes.add_textbox(_left(element.x, Inches), _top(element.y, Inches), _width(element.w, Inches), _height(element.h, Inches))
    _set_shape_name(shape, element)
    if element.style.fill:
        _set_fill_and_line(shape, element, deck, RGBColor)
    text_frame = shape.text_frame
    text_frame.clear()
    paragraphs = str(element.text or "").splitlines() or [""]
    for index, line in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = _paragraph_align(element.style.align, PP_ALIGN)
        for run in paragraph.runs:
            run.font.name = element.style.font_family or deck.theme.font_family
            run.font.size = Pt(max(1, int(element.style.font_size or 18)))
            run.font.bold = bool(element.style.bold)
            run.font.italic = bool(element.style.italic)
            run.font.underline = bool(element.style.underline)
            run.font.color.rgb = _rgb(element.style.color or deck.theme.ink, RGBColor)
    return shape


def _add_shape(slide, element: SlideElement, deck: DeckSpec, RGBColor, MSO_SHAPE, Inches):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if float(element.style.radius or 0) > 0 else MSO_SHAPE.RECTANGLE,
        _left(element.x, Inches),
        _top(element.y, Inches),
        _width(element.w, Inches),
        _height(element.h, Inches),
    )
    _set_shape_name(shape, element)
    _set_fill_and_line(shape, element, deck, RGBColor)
    if element.name:
        shape.text = element.name
    return shape


def _table_cells(element: SlideElement) -> list[list[str]]:
    rows = max(1, int(element.metadata.get("rows", 3) or 3))
    cols = max(1, int(element.metadata.get("cols", 3) or 3))
    raw = element.metadata.get("cells")
    cells: list[list[str]] = []
    if isinstance(raw, list):
        for row in raw[:rows]:
            cells.append([str(cell) for cell in row[:cols]] if isinstance(row, list) else [])
    while len(cells) < rows:
        cells.append([])
    for row_index, row in enumerate(cells):
        while len(row) < cols:
            row.append(f"Cell {row_index + 1}-{len(row) + 1}")
    return cells


def _add_table(slide, element: SlideElement, deck: DeckSpec, RGBColor, Inches, Pt):
    cells = _table_cells(element)
    rows = len(cells)
    cols = max(1, len(cells[0]))
    shape = slide.shapes.add_table(rows, cols, _left(element.x, Inches), _top(element.y, Inches), _width(element.w, Inches), _height(element.h, Inches))
    _set_shape_name(shape, element)
    table = shape.table
    header = bool(element.metadata.get("header", True))
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = format_formula_value(cells[r][c], cells=cells)
            try:
                cell.fill.solid()
                fill = element.metadata.get("header_fill") if header and r == 0 else element.metadata.get("body_fill")
                cell.fill.fore_color.rgb = _rgb(str(fill or element.style.fill or "#FFFFFF"), RGBColor)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = element.style.font_family or deck.theme.font_family
                        run.font.size = Pt(max(6, int(element.style.font_size or 12)))
                        run.font.color.rgb = _rgb(element.style.color or deck.theme.ink, RGBColor)
                        run.font.bold = bool(header and r == 0)
            except Exception:
                pass
    return shape


def _add_line(slide, element: SlideElement, deck: DeckSpec, RGBColor, MSO_CONNECTOR, Inches):
    x1 = _left(element.x, Inches)
    y1 = _top(element.y + element.h * 0.5, Inches)
    x2 = _left(element.x + element.w, Inches)
    y2 = y1
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    _set_shape_name(shape, element)
    try:
        shape.line.color.rgb = _rgb(element.style.stroke or element.style.color or deck.theme.accent, RGBColor)
        shape.line.width = int(max(1, float(element.style.stroke_width or 2)) * 12700)
    except Exception:
        pass
    return shape


def _add_image(slide, element: SlideElement, MSO_SHAPE, RGBColor, Inches):
    raw_source = str(element.source_path or "").strip()
    path = Path(raw_source) if raw_source else None
    if path is not None and path.is_file():
        shape = slide.shapes.add_picture(str(path), _left(element.x, Inches), _top(element.y, Inches), _width(element.w, Inches), _height(element.h, Inches))
        _set_shape_name(shape, element)
        return shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _left(element.x, Inches), _top(element.y, Inches), _width(element.w, Inches), _height(element.h, Inches))
    _set_shape_name(shape, element)
    shape.text = element.name or "Missing image"
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb("#F3F6FA", RGBColor)
        shape.line.color.rgb = _rgb("#2F6FED", RGBColor)
    except Exception:
        pass
    return shape


def _actor_poster_path(element: SlideElement) -> Path | None:
    for key in ("poster_path", "thumbnail_path", "preview_path", "render_path"):
        raw = str(element.metadata.get(key) or "").strip()
        if not raw:
            continue
        path = Path(raw)
        if path.is_file():
            return path
    return None


def _add_actor_placeholder(slide, element: SlideElement, deck: DeckSpec, MSO_SHAPE, RGBColor, Inches):
    poster = _actor_poster_path(element)
    if poster is not None:
        proxy = SlideElement.image(
            element.id,
            poster,
            x=element.x,
            y=element.y,
            w=element.w,
            h=element.h,
            kind="image",
            name=element.name,
        )
        proxy.opacity = element.opacity
        shape = _add_image(slide, proxy, MSO_SHAPE, RGBColor, Inches)
        try:
            shape.alt_text = f"{element.kind}: {element.source_path or element.name}"
        except Exception:
            pass
        return shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _left(element.x, Inches),
        _top(element.y, Inches),
        _width(element.w, Inches),
        _height(element.h, Inches),
    )
    _set_shape_name(shape, element)
    _set_fill_and_line(shape, element, deck, RGBColor)
    source_name = Path(str(element.source_path or "")).name
    title = element.name or element.kind.replace("_", " ").title()
    shape.text = f"{title}\n{element.kind.replace('_', ' ').title()}" + (f"\n{source_name}" if source_name else "")
    return shape


def _chart_type(value: str | None, XL_CHART_TYPE):
    normalized = str(value or "").strip().lower().replace("-", "_")
    if normalized in {"line", "lines"}:
        return XL_CHART_TYPE.LINE_MARKERS
    if normalized in {"pie", "donut", "doughnut"}:
        return XL_CHART_TYPE.PIE
    if normalized in {"horizontal_bar", "bar_horizontal"}:
        return XL_CHART_TYPE.BAR_CLUSTERED
    return XL_CHART_TYPE.COLUMN_CLUSTERED


def _chart_categories_and_values(element: SlideElement) -> tuple[list[str], list[float]]:
    raw_labels = element.metadata.get("labels") or ["A", "B", "C", "D"]
    raw_values = element.metadata.get("values") or [32, 58, 44, 72]
    labels = [str(label) for label in raw_labels] if isinstance(raw_labels, list) else ["A", "B", "C", "D"]
    source_values = list(raw_values) if isinstance(raw_values, list) else [32.0, 58.0, 44.0, 72.0]
    cells = [[labels[index] if index < len(labels) else f"Item {index + 1}", value] for index, value in enumerate(source_values)]
    count = max(1, min(len(labels), len(source_values)))
    out_labels = labels[:count] or ["A"]
    out_values: list[float] = []
    for value in source_values[:count]:
        try:
            out_values.append(float(evaluate_numeric_formula(value, cells=cells)))
        except Exception:
            out_values.append(0.0)
    return out_labels, out_values or [0.0]


def _add_chart(slide, element: SlideElement, deck: DeckSpec, ChartData, XL_CHART_TYPE, XL_LEGEND_POSITION, Inches, Pt):
    labels, values = _chart_categories_and_values(element)
    data = ChartData()
    data.categories = labels
    data.add_series(str(element.metadata.get("series_name") or element.name or "Series 1"), values)
    frame = slide.shapes.add_chart(
        _chart_type(str(element.metadata.get("chart_type") or "bar"), XL_CHART_TYPE),
        _left(element.x, Inches),
        _top(element.y, Inches),
        _width(element.w, Inches),
        _height(element.h, Inches),
        data,
    )
    _set_shape_name(frame, element)
    chart = frame.chart
    chart.has_legend = bool(element.metadata.get("show_legend", False))
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    title = str(element.metadata.get("title") or element.name or "").strip()
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
    try:
        chart.category_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.has_major_gridlines = True
    except Exception:
        pass
    try:
        chart.plots[0].vary_by_categories = True
    except Exception:
        pass
    return frame


def _shape_id_by_element_id(slide_xml: str) -> dict[str, int]:
    root = ET.fromstring(slide_xml.encode("utf-8"))
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    mapping: dict[str, int] = {}
    for node in root.findall(".//p:cNvPr", ns):
        name = str(node.attrib.get("name") or "")
        if not name:
            continue
        try:
            mapping[name] = int(node.attrib.get("id") or 0)
        except Exception:
            continue
    return mapping


def patch_pptx_animations(path: str | Path, deck: DeckSpec) -> Path:
    target = Path(path)
    temp = target.with_suffix(target.suffix + ".tmp")
    with zipfile.ZipFile(target, "r") as zin, zipfile.ZipFile(temp, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename.startswith("ppt/slides/slide") and info.filename.endswith(".xml"):
                try:
                    slide_index = int(Path(info.filename).stem.replace("slide", "")) - 1
                except Exception:
                    slide_index = -1
                if 0 <= slide_index < len(deck.slides):
                    xml = data.decode("utf-8")
                    slide = deck.slides[slide_index]
                    shape_ids = _shape_id_by_element_id(xml)
                    animations: list[str] = []
                    animated_shape_ids: list[int] = []
                    for element in sorted(slide.elements, key=animation_sequence_sort_key):
                        if not animation_is_active(element.animation):
                            continue
                        shape_id = shape_ids.get(element.id)
                        if not shape_id:
                            continue
                        animations.append(_animation_effect_xml(len(animations) + 1, shape_id, element))
                        animated_shape_ids.append(shape_id)
                    if animations:
                        timing = _timing_xml(animations, animated_shape_ids)
                        if "<p:timing>" not in xml:
                            marker = "</p:clrMapOvr>"
                            if marker in xml:
                                xml = xml.replace(marker, marker + timing, 1)
                            else:
                                xml = xml.replace("</p:sld>", timing + "</p:sld>", 1)
                        data = xml.encode("utf-8")
            zout.writestr(info, data)
    shutil.move(str(temp), str(target))
    return target


def write_pptx_compatible(deck: DeckSpec, path: str | Path, *, include_animations: bool = True) -> Path:
    Presentation, ChartData, RGBColor, MSO_CONNECTOR, MSO_SHAPE, PP_ALIGN, XL_CHART_TYPE, XL_LEGEND_POSITION, Inches, Pt = _require_pptx()
    target = Path(path)
    target.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    for slide_spec in deck.slides:
        slide = prs.slides.add_slide(blank)
        try:
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = _rgb(slide_spec.background or deck.theme.background, RGBColor)
        except Exception:
            pass
        for element in sorted(slide_spec.elements, key=lambda row: int(row.z_index)):
            if not element.visible:
                continue
            if element.kind in {"text", "typography_actor"}:
                _add_text(slide, element, deck, RGBColor, PP_ALIGN, Inches, Pt)
            elif element.kind in {"image", "timeline_moment", "screen_capture"}:
                _add_image(slide, element, MSO_SHAPE, RGBColor, Inches)
            elif element.kind == "table":
                _add_table(slide, element, deck, RGBColor, Inches, Pt)
            elif element.kind == "chart":
                _add_chart(slide, element, deck, ChartData, XL_CHART_TYPE, XL_LEGEND_POSITION, Inches, Pt)
            elif element.kind == "line":
                _add_line(slide, element, deck, RGBColor, MSO_CONNECTOR, Inches)
            elif element.kind in ACTOR_PLACEHOLDER_KINDS:
                _add_actor_placeholder(slide, element, deck, MSO_SHAPE, RGBColor, Inches)
            else:
                _add_shape(slide, element, deck, RGBColor, MSO_SHAPE, Inches)
    prs.save(target)
    if include_animations:
        patch_pptx_animations(target, deck)
    return target


__all__ = ["patch_pptx_animations", "write_pptx_compatible"]
