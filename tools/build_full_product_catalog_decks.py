from __future__ import annotations

import hashlib
import json
import shutil
import sys
import warnings
import zipfile
from dataclasses import dataclass
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageStat

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))


WORKSPACE = ROOT.parent / "ReviewAutomationWorkspace"
TEMPLATES = WORKSPACE / "source_assets" / "templates"
TMP = WORKSPACE / "tmp"
OUT = WORKSPACE / "outputs" / "product_catalog_full"
FRESH_CAPTURE_ROOT = TMP / "fresh_review_recapture"
STRICT_REPORT = FRESH_CAPTURE_ROOT / "missing_strict_catalog_assets.txt"
BUILD_BLOCKED_REPORT = OUT / "build_blocked_report.md"
BUILD_WARNINGS: list[str] = []
FORBIDDEN_FINAL_CAPTURE_MARKERS = (
    "fresh_first_slide_capture",
    "actual_3d_viewer_capture",
    "debugcapture",
)
SLIDE_W = 1672
SLIDE_H = 941
EDITOR_VIEWER_BOX = (206, 126, 890, 503)
EDITOR_REFERENCE_SIZE = (1480, 920)
EDITOR_VIEWER_BOX_BY_ASSET = {
    # The color-grading workspace opens a large right workbench. The generic
    # viewer box includes part of that workbench, so it can miss a black viewer.
    "color_before_after_editor": (206, 126, 592, 503),
}
AR_PBR_MATCHED_EDITOR_CAPTURE = (
    FRESH_CAPTURE_ROOT
    / "ar_pbr_statue_composite"
    / "editor_ar_pbr_statue_composite_action.png"
)
SPEC_CLOSING_BONSAI = TMP / "catalog_spec_closing" / "bonsai_blue_pot_cutout_v1.png"
SPEC_CLOSING_SHADOW_MODE = "pot_contact_only"
AR_PBR_SAME_ASSET_RULE = (
    "AR/PBR same-asset rule: capture the approved plaster statue/bust loaded "
    "into the editor video viewer, scaled large enough to read, before building "
    "the catalog. Do not fall back to Camera_01 or any older debug capture."
)
FEATURE_EVIDENCE_RULE = (
    "Feature-evidence lock: this slide needs a fresh, feature-specific editor "
    "capture whose visible UI matches the slide claim. Do not substitute a "
    "generic editor screenshot, standalone viewer, raw video frame, or older "
    "debug capture."
)
MULTI_MONITOR_OVERVIEW_RULE = (
    "Multi-monitor overview lock: slide 1 must prove the documented monitor "
    "roles. Center is the only monitor with the main video preview/frame viewer "
    "and must show the Lamborghini edit, a long multi-track timeline, and AI "
    "command/chat. Left is actor/asset support with Live2D plus real AR/PBR/3D "
    "and MMD/VRM-style support surfaces. Right is node-dominant with sound/audio "
    "workbench evidence. A semantic sidecar may not be auto-stamped from an "
    "existing image; it must record these role-specific facts from a real "
    "TigerCapture capture."
)
COMPARE_EVIDENCE_RULE = (
    "Before/after evidence lock: apply the grade/filter/node through the editor "
    "action surface, enable ui.viewer.compare.set(split or before), capture the "
    "real changed editor state, and validate that the viewer is not black. The "
    "capture must also write a sidecar contract proving non-neutral parameter "
    "changes and a visible delta; original/neutral-looking output is invalid."
)
COLOR_IPAD_DETAIL_RULE = (
    "Color iPad detail lock: capture only the color-control detail surface "
    "(wheels, curves, scopes, tone controls, and sliders). The iPad source must "
    "not contain the video viewer, media pool, or timeline."
)
ACTOR_COMPOSITE_RULE = (
    "Actor-composite evidence lock: the actor must be visible inside the editor "
    "video viewer and on the timeline/actor lane. A standalone actor viewer is "
    "not enough for the product-catalog slide."
)
VTUBER_STUDIO_RULE = (
    "VTuber Studio evidence lock: capture the real Program Output / Source "
    "Tracking / Avatar Mapping / Studio Controls layout. The laptop/main frame "
    "must be the full VTuber Studio workspace. The iPad/detail frame must be "
    "Program Output only, not Source Tracking, Avatar Mapping, a duplicated "
    "workspace, or a generic editor crop. When using the Trump Performance "
    "Source, Avatar Mapping and Program Output must show upper-body avatar "
    "evidence: head, neck, shoulders, and upper torso. A face-only VRM meta "
    "thumbnail is invalid. The performance source must not be shown as the "
    "final Program Output. Product evidence must use the VTuber VRM GPU renderer "
    "(`vrm_mtoon_gpu`); software VRM fallback output is invalid because it can "
    "render dotted/point-cloud avatars."
)
VTUBER_VALIDATED_ASSETS = {
    "vtuber_studio_editor",
    "vtuber_studio_program_output",
    "vtuber_studio_tracking_mapping",
    "vtuber_studio_avatar_mapping",
}
VTUBER_FORBIDDEN_EVIDENCE_MARKERS = (
    "polyhaven_pbr_camera_scene",
    "camera_01",
    "ar_pbr_camera",
    "ar_pbr_nexus",
    "ar_pbr_statue",
    "marmoset",
    "full-gpu",
    "full_gpu",
    "debug proof",
)
VTUBER_REQUIRED_INPUT_MARKERS = {
    "trump_performance_source": "trump_oval_office_live_gnzweo_hfe0.mp4",
    "program_output_background": "south korea 4k drone video",
    "vrm_avatar_target": "milica_v1.3.vrm",
}
VIEWER_VALIDATED_ASSETS = {
    "main_editor_current",
    "media_pool_current",
    "timeline_current_editor",
    "overview_center_editor",
    "effects_before_after_editor",
    "transitions_editor",
    "typography_editor",
    "keyframe_motion_editor",
    "color_before_after_editor",
    "node_before_after_editor",
    "node_effect_before_after_editor",
    "live2d_composite_editor",
    "mmd_composite_editor",
    "ar_statue_editor",
    "export_editor_current",
}
COLOR_IPAD_ALLOWED_NAME_MARKERS = (
    "color_controls",
    "color_wheels",
    "color_curves",
    "color_scopes",
    "grading_controls",
    "grade_controls",
    "color_panel",
)
DETAIL_IPAD_FORBIDDEN_NAME_MARKERS = (
    "editor_",
    "full_editor",
    "timeline",
    "viewer",
    "media_pool",
    "overview",
)
DETAIL_IPAD_FULL_EDITOR_MIN_SIZE = (1100, 650)
DETAIL_IPAD_MIN_SIZE = (320, 220)
COMPARE_VALIDATED_ASSETS = {
    "effects_before_after_editor",
    "color_before_after_editor",
    "node_before_after_editor",
    "node_graph_actual",
    "node_effect_before_after_editor",
}
COMPARE_ALLOWED_MODES = {"split", "before", "wipe", "before_after", "original_after"}
COMPARE_ACTION_REQUIREMENTS: dict[str, tuple[str, ...]] = {
    "effects_before_after_editor": (
        "clip.set_filter",
        "effect.apply_to_clip",
        "effect.hover_or_select",
        "node.graph.set",
    ),
    "color_before_after_editor": (
        "clip.set_color_grade",
        "color.set_wheels_or_curves",
        "color.workspace.open",
        "color.grade.apply",
    ),
    "node_before_after_editor": (
        "node.graph.set",
        "node.add",
        "node.connect",
        "node.set_param",
        "clip.set_filter",
    ),
    "node_graph_actual": (
        "node.graph.set",
        "node.add",
        "node.connect",
        "node.set_param",
    ),
    "node_effect_before_after_editor": (
        "node.graph.set",
        "node.add",
        "node.connect",
        "node.set_param",
        "clip.set_filter",
    ),
}
COMPARE_SOURCE_REPORT_REQUIRED_CHECKS: dict[str, tuple[str, ...]] = {
    "color_before_after_editor": (
        "viewer_frame_visible",
        "color_dock_viewer_reforced",
        "viewer_compare_split",
    ),
}
COMPARE_ACTION_KEYS = (
    "executed_actions",
    "action_sequence",
    "actions",
    "steps",
    "action_log",
)
SEMANTIC_CAPTURE_CONTRACTS: dict[str, dict[str, object]] = {
    "overview_left_workspace": {
        "contract": "multi_monitor_left_workspace_v1",
        "contains_all": ("live2d_viewer", "ar_pbr_viewer"),
        "forbidden": (
            "main_video_preview",
            "timeline_main_video",
            "center_editor_duplicate",
            "video_frame_viewer",
        ),
    },
    "overview_center_editor": {
        "contract": "multi_monitor_center_editor_v1",
        "contains_all": ("main_video_preview", "timeline", "ai_command"),
        "forbidden": (
            "node_graph",
            "left_workspace_duplicate",
            "right_workspace_duplicate",
            "side_monitor_workspace",
        ),
    },
    "overview_right_workspace": {
        "contract": "multi_monitor_right_workspace_v1",
        "contains_all": ("node_graph",),
        "contains_any": ("sound_editor", "audio_workbench", "audio_visualizer"),
        "forbidden": ("main_video_preview", "center_editor_duplicate", "video_frame_viewer"),
    },
    "effects_hover_detail": {
        "contract": "effect_detail_v1",
        "contains_any": ("effect_controls", "effect_preview", "effect_library", "effect_before_after"),
        "forbidden": ("timeline_only", "generic_editor_crop"),
    },
    "transition_detail": {
        "contract": "transition_detail_v1",
        "contains_any": ("transition_controls", "transition_preview", "transition_handle"),
        "forbidden": ("timeline_only", "generic_editor_crop"),
    },
    "typography_detail": {
        "contract": "typography_detail_v1",
        "contains_all": ("typography_controls", "multiple_text_styles"),
        "forbidden": ("timeline_only", "generic_editor_crop", "single_tiny_caption"),
    },
    "keyframe_detail": {
        "contract": "keyframe_detail_v1",
        "contains_any": ("keyframe_controls", "curve_editor", "transform_keyframes", "opacity_keyframes"),
        "forbidden": ("timeline_only", "generic_editor_crop"),
    },
    "ppt_maker_detail": {
        "contract": "ppt_maker_detail_v1",
        "contains_any": ("ppt_actions", "element_inspector", "export_snapshot", "validation_contact_sheet"),
        "forbidden": ("full_editor", "generic_powerpoint_mockup", "debugcapture_source"),
    },
    "node_graph_actual": {
        "contract": "node_graph_detail_v1",
        "contains_all": ("node_graph",),
        "contains_any": ("selected_node_params", "node_controls", "before_after_node_result"),
        "forbidden": ("black_viewer", "generic_editor_crop"),
    },
    "node_effect_library_detail": {
        "contract": "node_effect_library_detail_v1",
        "contains_any": ("node_effect_library", "effect_node_controls", "before_after_node_result"),
        "forbidden": ("timeline_only", "generic_editor_crop"),
    },
    "live2d_composite_editor": {
        "contract": "live2d_composite_editor_v1",
        "contains_all": ("live2d_actor", "actor_lane"),
        "forbidden": ("raw_video_only", "standalone_viewer_only"),
    },
    "live2d_actor_detail": {
        "contract": "live2d_viewer_detail_v1",
        "contains_all": ("live2d_viewer", "live2d_actor"),
        "forbidden": ("generic_editor_crop", "raw_video_only"),
    },
    "mmd_composite_editor": {
        "contract": "mmd_composite_editor_v1",
        "contains_all": ("mmd_character", "actor_lane"),
        "forbidden": ("live2d_substitute", "raw_video_only", "standalone_viewer_only"),
    },
    "mmd_character_detail": {
        "contract": "mmd_viewer_detail_v1",
        "contains_all": ("mmd_viewer", "mmd_character"),
        "forbidden": ("live2d_substitute", "generic_editor_crop", "raw_video_only"),
    },
}
DETAIL_SEMANTIC_ASSETS = {
    "effects_hover_detail",
    "transition_detail",
    "typography_detail",
    "keyframe_detail",
    "ppt_maker_detail",
    "node_graph_actual",
    "node_effect_library_detail",
    "live2d_actor_detail",
    "mmd_character_detail",
}


def _contract_path_for_capture(path: Path) -> Path:
    return path.with_suffix(".capture-contract.json")


def _is_non_neutral_changed_param(value: object) -> bool:
    if isinstance(value, dict):
        after = value.get("after", value.get("value", value.get("applied")))
        neutral = value.get("neutral", value.get("default", value.get("before")))
        if after is None:
            return bool(value.get("non_neutral") or value.get("changed"))
        if isinstance(after, (int, float)) and isinstance(neutral, (int, float)):
            return abs(float(after) - float(neutral)) > 1e-6
        return str(after) != str(neutral)
    if isinstance(value, str):
        return bool(value.strip())
    return value not in (None, False, 0, 0.0)


def _action_names_from_value(value: object) -> list[str]:
    names: list[str] = []
    if isinstance(value, str):
        if value.strip():
            names.append(value.strip())
        return names
    if isinstance(value, dict):
        for key in ("action", "id", "name", "action_id", "command"):
            item = value.get(key)
            if isinstance(item, str) and item.strip():
                names.append(item.strip())
        for key in COMPARE_ACTION_KEYS:
            names.extend(_action_names_from_value(value.get(key)))
        return names
    if isinstance(value, (list, tuple, set)):
        for item in value:
            names.extend(_action_names_from_value(item))
    return names


def _resolve_contract_sidecar_path(contract_path: Path, value: object) -> Path | None:
    raw = str(value or "").strip()
    if not raw:
        return None
    path = Path(raw)
    if not path.is_absolute():
        path = contract_path.parent / path
    return path


def _compare_action_evidence(
    data: dict[str, object],
    contract_path: Path,
) -> tuple[bool, str, list[str]]:
    names: list[str] = []
    for key in COMPARE_ACTION_KEYS:
        names.extend(_action_names_from_value(data.get(key)))
    action_surface = str(data.get("action_surface") or "").strip()
    if action_surface:
        names.extend(
            token.strip()
            for token in action_surface.replace("+", ",").replace("->", ",").split(",")
            if token.strip()
        )

    source_report_path = _resolve_contract_sidecar_path(contract_path, data.get("source_report"))
    if source_report_path is not None:
        if not source_report_path.exists():
            return False, f"Compare contract source_report is missing: {source_report_path}", names
        try:
            report = json.loads(source_report_path.read_text(encoding="utf-8"))
        except Exception as exc:
            return False, f"Could not read compare source_report {source_report_path}: {exc}", names
        if report.get("ok") is False:
            return False, f"Compare source_report is not successful: {source_report_path}", names
        for key in COMPARE_ACTION_KEYS:
            names.extend(_action_names_from_value(report.get(key)))
    elif not any(data.get(key) for key in COMPARE_ACTION_KEYS) and not bool(
        data.get("real_action_report_embedded")
    ):
        return (
            False,
            "Compare contract has no source_report or embedded action execution log.",
            names,
        )

    normalized = [name.casefold() for name in names if str(name).strip()]
    return True, "compare action evidence ok", normalized


def _compare_source_report_checks_are_ready(
    name: str,
    data: dict[str, object],
    contract_path: Path,
) -> tuple[bool, str]:
    required = COMPARE_SOURCE_REPORT_REQUIRED_CHECKS.get(name, ())
    if not required:
        return True, "no compare source_report checks required"
    source_report_path = _resolve_contract_sidecar_path(contract_path, data.get("source_report"))
    if source_report_path is None:
        return False, f"Before/after contract for {name} has no source_report for required capture checks."
    if not source_report_path.exists():
        return False, f"Compare contract source_report is missing: {source_report_path}"
    try:
        report = json.loads(source_report_path.read_text(encoding="utf-8"))
    except Exception as exc:
        return False, f"Could not read compare source_report {source_report_path}: {exc}"
    checks = report.get("checks") if isinstance(report.get("checks"), dict) else {}
    missing_or_false = [key for key in required if not bool(checks.get(key))]
    if missing_or_false:
        return (
            False,
            f"Before/after contract for {name} has failed capture checks in source_report: "
            + ", ".join(missing_or_false),
        )
    return True, "compare source_report checks ok"


def _contains_action(actions: list[str], *needles: str) -> bool:
    for needle in needles:
        needle_cf = needle.casefold()
        if any(needle_cf in action for action in actions):
            return True
    return False


def _compare_capture_contract_is_ready(name: str, image_path: Path) -> tuple[bool, str]:
    contract_path = _contract_path_for_capture(image_path)
    if not contract_path.exists():
        return (
            False,
            "Missing before/after capture contract sidecar. Recapture this page "
            f"and write {contract_path.name} with changed_params, compare mode, "
            "non-neutral confirmation, and visible_delta=true.",
        )
    try:
        data = json.loads(contract_path.read_text(encoding="utf-8"))
    except Exception as exc:
        return False, f"Could not read before/after capture contract: {exc}"

    mode = str(data.get("viewer_compare_mode") or data.get("compare_mode") or "").strip().lower()
    if mode not in COMPARE_ALLOWED_MODES:
        return False, f"Before/after contract for {name} has invalid compare mode: {mode!r}"
    if bool(data.get("neutral_identity")) or bool(data.get("result_matches_original")):
        return False, f"Before/after contract for {name} says the result is neutral/original-like."
    if not bool(data.get("visible_delta")):
        return False, f"Before/after contract for {name} does not confirm a visible image delta."

    changed = data.get("changed_params") or data.get("parameters_changed") or data.get("applied_params")
    if isinstance(changed, dict):
        changed_values = list(changed.values())
    elif isinstance(changed, list):
        changed_values = changed
    else:
        changed_values = []
    if not changed_values:
        return False, f"Before/after contract for {name} has no changed_params."
    if not bool(data.get("non_neutral_params_confirmed")) and not any(
        _is_non_neutral_changed_param(item) for item in changed_values
    ):
        return False, f"Before/after contract for {name} does not prove non-neutral parameter values."
    if bool(data.get("preset_values_unknown")) and not str(
        data.get("preset_source") or data.get("preset_reference") or data.get("preset_reference_url") or ""
    ).strip():
        return False, (
            f"Before/after contract for {name} used unknown preset values but "
            "does not record an internet/reference source."
        )
    ok, reason, actions = _compare_action_evidence(data, contract_path)
    if not ok:
        return False, reason
    ok, reason = _compare_source_report_checks_are_ready(name, data, contract_path)
    if not ok:
        return False, reason
    if not (
        _contains_action(actions, "ui.viewer.compare.set")
        or bool(data.get("compare_action_executed"))
        or bool(data.get("viewer_compare_action_executed"))
    ):
        return (
            False,
            f"Before/after contract for {name} does not prove ui.viewer.compare.set was executed.",
        )
    required_actions = COMPARE_ACTION_REQUIREMENTS.get(name, ())
    if required_actions and not _contains_action(actions, *required_actions):
        return (
            False,
            f"Before/after contract for {name} does not prove the required feature action was executed: "
            + ", ".join(required_actions),
        )
    return True, f"before/after contract ok: {contract_path.name}"


def _contract_tags(data: dict[str, object]) -> set[str]:
    tags: set[str] = set()
    for key in (
        "contains",
        "visible_features",
        "evidence_tags",
        "role_tags",
        "forbidden_tags",
    ):
        value = data.get(key)
        if isinstance(value, str):
            tags.add(value.strip().casefold())
        elif isinstance(value, list):
            tags.update(str(item).strip().casefold() for item in value if str(item).strip())
    for key, value in data.items():
        if isinstance(value, bool) and value:
            tags.add(str(key).strip().casefold())
    return tags


def _contract_feature_present(data: dict[str, object], tags: set[str], *names: str) -> bool:
    falsey_strings = {"", "0", "false", "no", "none", "null", "missing", "n/a"}
    for name in names:
        key = str(name).strip()
        if not key:
            continue
        if key.casefold() in tags:
            return True
        value = data.get(key)
        if isinstance(value, bool):
            if value:
                return True
        elif isinstance(value, (int, float)):
            if value:
                return True
        elif isinstance(value, str):
            if value.strip().casefold() not in falsey_strings:
                return True
        elif isinstance(value, (list, tuple, set)):
            if value:
                return True
    return False


def _contract_json_text(data: dict[str, object]) -> str:
    try:
        return json.dumps(data, ensure_ascii=False, sort_keys=True).casefold()
    except Exception:
        return str(data).casefold()


def _contract_text_contains(data: dict[str, object], *tokens: str) -> bool:
    text = _contract_json_text(data)
    return any(str(token).casefold() in text for token in tokens)


def _contract_numeric_at_least(data: dict[str, object], minimum: float, *names: str) -> bool:
    for name in names:
        value = _to_float(data.get(name))
        if value is not None and value >= minimum:
            return True
    return False


def _overview_contract_is_ready(name: str, data: dict[str, object], tags: set[str]) -> tuple[bool, str]:
    if name not in {"overview_left_workspace", "overview_center_editor", "overview_right_workspace"}:
        return True, "not a multi-monitor overview contract"

    role = str(data.get("monitor_role") or data.get("screen_role") or data.get("role") or "").casefold()
    expected_role = {
        "overview_left_workspace": "left",
        "overview_center_editor": "center",
        "overview_right_workspace": "right",
    }[name]
    if expected_role not in role and f"{expected_role}_monitor" not in tags:
        return (
            False,
            f"Multi-monitor overview contract for {name} must identify monitor_role={expected_role!r}.",
        )

    if not _contract_feature_present(
        data,
        tags,
        "real_tigercapture_capture",
        "actual_tigercapture_window",
        "actual_window_capture",
        "ui_action_capture",
        "review_action_capture",
    ):
        return (
            False,
            f"Multi-monitor overview contract for {name} must prove it came from a real TigerCapture window/action capture.",
        )

    if name == "overview_center_editor":
        if not (
            _contract_feature_present(data, tags, "lamborghini_clip", "center_media_lamborghini", "lamborghini_video")
            or _contract_text_contains(data, "lamborghini")
        ):
            return (
                False,
                "Center monitor overview must use the Lamborghini edit media called out by the multi-monitor rules.",
            )
        if not _contract_feature_present(data, tags, "long_timeline", "timeline_long_enough"):
            return False, "Center monitor overview must prove a long timeline, not a short/simple edit."
        if not (
            _contract_feature_present(data, tags, "multi_track_timeline", "multiple_tracks_visible")
            or _contract_numeric_at_least(data, 2, "visible_track_count", "timeline_track_count")
        ):
            return (
                False,
                "Center monitor overview must prove a multi-track timeline with at least two visible tracks.",
            )
        if not _contract_feature_present(data, tags, "ai_command_secondary", "ai_chat_visible", "local_ai_visible"):
            return (
                False,
                "Center monitor overview must show AI command/chat as a secondary center workspace element.",
            )
        if _contract_feature_present(
            data,
            tags,
            "macro_eye",
            "macro_face",
            "human_closeup",
            "face_closeup",
            "body_closeup",
            "person_closeup",
        ):
            return False, "Center monitor overview may not use macro eye/face/body close-up evidence."
        return True, "multi-monitor center overview contract ok"

    if name == "overview_left_workspace":
        if not _contract_feature_present(
            data,
            tags,
            "mmd_viewer",
            "mmd_character",
            "vrm_studio",
            "vtuber_studio",
            "avatar_mapping",
            "vrm_avatar",
            "actor_support_surface",
        ):
            return (
                False,
                "Left monitor overview must include MMD/VRM/VTuber or avatar-mapping support, not only Live2D/3D.",
            )
        if not _contract_feature_present(data, tags, "asset_preset_support", "asset_browser", "preset_support"):
            return False, "Left monitor overview must include actor/asset or preset support surfaces."
        if not _contract_feature_present(data, tags, "cubemap_hidden", "neutral_3d_background", "ar_pbr_background_hidden"):
            return (
                False,
                "Left monitor AR/PBR evidence must prove the 3D viewer background/cubemap is hidden or neutral.",
            )
        return True, "multi-monitor left overview contract ok"

    if name == "overview_right_workspace":
        if not _contract_feature_present(data, tags, "node_graph_dominant", "large_node_graph", "node_workspace_primary"):
            return False, "Right monitor overview must prove the node graph is the dominant workspace."
        if not _contract_feature_present(
            data,
            tags,
            "sound_editor",
            "audio_workbench",
            "audio_visualizer",
            "sound_secondary",
            "audio_mixer",
            "audio_scopes",
        ):
            return False, "Right monitor overview must include sound/audio workbench evidence."
        if _contract_feature_present(data, tags, "generic_workbench_inspector", "fake_node_graph"):
            return False, "Right monitor overview cannot be a generic inspector or fake node graph."
        return True, "multi-monitor right overview contract ok"

    return True, "multi-monitor overview contract ok"


def _semantic_capture_contract_is_ready(name: str, image_path: Path) -> tuple[bool, str]:
    rules = SEMANTIC_CAPTURE_CONTRACTS.get(name)
    if not rules:
        return True, "no semantic capture contract required"
    contract_path = _contract_path_for_capture(image_path)
    if not contract_path.exists():
        return (
            False,
            "Missing semantic capture contract sidecar. Recapture this feature "
            f"and write {contract_path.name}; do not reuse a similar-looking "
            "editor screenshot from another page.",
        )
    try:
        data = json.loads(contract_path.read_text(encoding="utf-8"))
    except Exception as exc:
        return False, f"Could not read semantic capture contract: {exc}"

    expected = str(rules.get("contract") or "").strip()
    actual = str(
        data.get("semantic_contract")
        or data.get("catalog_contract")
        or data.get("contract")
        or ""
    ).strip()
    if actual != expected:
        return False, f"Semantic contract mismatch for {name}: expected {expected!r}, got {actual!r}"

    tags = _contract_tags(data)
    required = {str(item).casefold() for item in rules.get("contains_all", ())}
    missing = sorted(required - tags)
    if missing:
        return False, f"Semantic contract for {name} is missing required visible features: {', '.join(missing)}"

    options = {str(item).casefold() for item in rules.get("contains_any", ())}
    if options and not (options & tags):
        return (
            False,
            f"Semantic contract for {name} must include at least one of: "
            + ", ".join(sorted(options)),
        )

    forbidden = {str(item).casefold() for item in rules.get("forbidden", ())}
    present_forbidden = sorted(forbidden & tags)
    if present_forbidden:
        return False, f"Semantic contract for {name} contains forbidden evidence tags: {', '.join(present_forbidden)}"
    if bool(data.get("substituted_from_other_feature")):
        return False, f"Semantic contract for {name} admits cross-feature substitution."
    ok, reason = _overview_contract_is_ready(name, data, tags)
    if not ok:
        return False, reason
    if name in {"live2d_composite_editor", "mmd_composite_editor"} and not bool(
        data.get("main_viewer_actor_visible")
        or data.get("actor_visible_in_main_viewer")
        or data.get("viewer_actor_overlay_visible")
    ):
        return (
            False,
            f"Semantic contract for {name} does not prove the actor is visible inside the main editor viewer.",
        )
    if name in {"mmd_composite_editor", "mmd_character_detail"}:
        ok, reason = _mmd_motion_frame_contract_is_ready(name, data)
        if not ok:
            return False, reason
    if name == "typography_detail":
        try:
            layer_count = int(
                data.get("visible_text_layer_count")
                or data.get("text_layer_count")
                or data.get("typography_actor_count")
                or 0
            )
        except Exception:
            layer_count = 0
        if layer_count < 4:
            return (
                False,
                "Semantic contract for typography_detail must prove at least "
                "four visible typography layers: headline, secondary line, "
                "multilingual sample, and smaller caption/body text.",
            )
        required_flags = {
            "large_headline_visible": "large headline",
            "secondary_text_visible": "secondary text layer",
            "multilingual_text_visible": "multilingual text sample",
            "small_caption_text_visible": "smaller caption/body text",
        }
        missing_flags = [label for key, label in required_flags.items() if not bool(data.get(key))]
        if missing_flags:
            return (
                False,
                "Semantic contract for typography_detail is missing visible typography proof: "
                + ", ".join(missing_flags),
            )
    return True, f"semantic contract ok: {contract_path.name}"


def _to_float(value: object) -> float | None:
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _mmd_motion_frame_contract_is_ready(name: str, data: dict[str, object]) -> tuple[bool, str]:
    if data.get("first_frame_used") is not False:
        return (
            False,
            f"Semantic contract for {name} must explicitly set first_frame_used=false. "
            "MMD catalog evidence may not use frame 0.",
        )

    position = str(
        data.get("capture_frame_position")
        or data.get("motion_frame_policy")
        or data.get("frame_selection")
        or ""
    ).casefold()
    progress = _to_float(data.get("capture_progress") or data.get("capture_normalized_time"))
    time_ms = _to_float(data.get("capture_time_ms") or data.get("mmd_capture_time_ms"))
    mid_position = any(token in position for token in ("middle", "mid", "motion", "active"))
    mid_progress = progress is not None and 0.2 <= progress <= 0.8
    positive_time = time_ms is not None and time_ms >= 1000.0
    if not (mid_position or mid_progress or positive_time):
        return (
            False,
            f"Semantic contract for {name} must prove a middle/active MMD frame "
            "with capture_frame_position='mid_motion', capture_progress 0.2..0.8, "
            "or capture_time_ms >= 1000.",
        )

    pose_delta = _to_float(data.get("motion_pose_delta") or data.get("mmd_motion_pose_delta"))
    active = bool(
        data.get("mmd_motion_active")
        or data.get("motion_activity_visible")
        or data.get("motion_controls_active")
        or (pose_delta is not None and pose_delta > 0.0)
    )
    if not active:
        return (
            False,
            f"Semantic contract for {name} must prove visible MMD motion activity, "
            "not a static first-frame or idle pose.",
        )
    return True, "MMD middle motion frame contract ok"


def _semantic_capture_visual_is_ready(name: str, image_path: Path) -> tuple[bool, str]:
    if name not in SEMANTIC_CAPTURE_CONTRACTS:
        return True, "no semantic visual contract required"
    try:
        img = Image.open(image_path).convert("RGB")
    except Exception as exc:
        return False, f"Could not inspect semantic capture pixels: {exc}"
    if img.width < 320 or img.height < 220:
        return False, f"Semantic capture is too small to be product evidence: {img.width}x{img.height}"

    gray = img.convert("L")
    stat = ImageStat.Stat(gray)
    mean = float(stat.mean[0])
    stddev = float(stat.stddev[0])
    if mean < 8.0 and stddev < 8.0:
        return False, f"Semantic capture is black/blank: mean_luma={mean:.2f}, stddev={stddev:.2f}"
    if stddev < 2.5:
        return False, f"Semantic capture is visually flat/empty: mean_luma={mean:.2f}, stddev={stddev:.2f}"
    if mean > 247.0 and stddev < 9.0:
        return False, f"Semantic capture is nearly blank white: mean_luma={mean:.2f}, stddev={stddev:.2f}"

    if name in DETAIL_SEMANTIC_ASSETS:
        small_w = min(360, img.width)
        small_h = min(240, img.height)
        small = img.resize((small_w, small_h), Image.Resampling.BILINEAR)
        corner = 16
        samples: list[tuple[int, int, int]] = []
        for box in (
            (0, 0, corner, corner),
            (small_w - corner, 0, small_w, corner),
            (0, small_h - corner, corner, small_h),
            (small_w - corner, small_h - corner, small_w, small_h),
        ):
            with warnings.catch_warnings():
                warnings.simplefilter("ignore", DeprecationWarning)
                samples.extend(list(small.crop(box).getdata()))
        bg = tuple(
            sorted(pixel[channel] for pixel in samples)[len(samples) // 2]
            for channel in range(3)
        )
        active_pixels = 0
        active_rows = 0
        row_stride = small_w
        with warnings.catch_warnings():
            warnings.simplefilter("ignore", DeprecationWarning)
            pixels = list(small.getdata())
        for y in range(small_h):
            row = pixels[y * row_stride : (y + 1) * row_stride]
            row_active = 0
            for r, g, b in row:
                if abs(r - bg[0]) + abs(g - bg[1]) + abs(b - bg[2]) > 34:
                    row_active += 1
            active_pixels += row_active
            if row_active / float(small_w) > 0.04:
                active_rows += 1
        active_pixel_ratio = active_pixels / float(small_w * small_h)
        active_row_ratio = active_rows / float(small_h)
        if active_pixel_ratio < 0.10:
            return (
                False,
                f"Detail capture has too little visible content: active_pixels={active_pixel_ratio:.3f}. "
                "Do not use blank panels or tiny PPT/timeline fragments.",
            )
        if active_row_ratio < 0.24:
            return (
                False,
                f"Detail capture is a thin strip, not a meaningful feature detail: active_rows={active_row_ratio:.3f}.",
            )
    return True, f"semantic pixels ok: mean_luma={mean:.2f}, stddev={stddev:.2f}"


def _image_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _image_dhash(img: Image.Image, *, size: int = 8) -> int:
    gray = img.convert("L").resize((size + 1, size), Image.Resampling.BILINEAR)
    pixels = list(gray.getdata())
    bits = 0
    for y in range(size):
        row = y * (size + 1)
        for x in range(size):
            bits = (bits << 1) | int(pixels[row + x] > pixels[row + x + 1])
    return bits


def _hamming_distance(left: int, right: int) -> int:
    return int(left ^ right).bit_count()


def _screen_crop_signature(img: Image.Image) -> str:
    normalized = img.convert("RGB").resize((160, 90), Image.Resampling.BILINEAR)
    return hashlib.sha256(normalized.tobytes()).hexdigest()


def _screen_region_is_catalog_ready(img: Image.Image, *, label: str) -> tuple[bool, str]:
    if img.width < 120 or img.height < 80:
        return False, f"{label} screen region is too small: {img.width}x{img.height}"
    rgb = img.convert("RGB")
    gray = rgb.convert("L")
    stat = ImageStat.Stat(gray)
    mean = float(stat.mean[0])
    stddev = float(stat.stddev[0])
    if mean < 6.0 and stddev < 6.0:
        return False, f"{label} screen region is black/blank: mean_luma={mean:.2f}, stddev={stddev:.2f}"
    if mean > 248.0 and stddev < 8.0:
        return False, f"{label} screen region is blank white: mean_luma={mean:.2f}, stddev={stddev:.2f}"
    if stddev < 4.5:
        return False, f"{label} screen region is visually flat/empty: mean_luma={mean:.2f}, stddev={stddev:.2f}"

    small_w = min(360, rgb.width)
    small_h = min(220, rgb.height)
    small = rgb.resize((small_w, small_h), Image.Resampling.BILINEAR)
    corner = max(8, min(18, small_w // 10, small_h // 10))
    samples: list[tuple[int, int, int]] = []
    with warnings.catch_warnings():
        warnings.simplefilter("ignore", DeprecationWarning)
        for box in (
            (0, 0, corner, corner),
            (small_w - corner, 0, small_w, corner),
            (0, small_h - corner, corner, small_h),
            (small_w - corner, small_h - corner, small_w, small_h),
        ):
            samples.extend(list(small.crop(box).getdata()))
    bg = tuple(
        sorted(pixel[channel] for pixel in samples)[len(samples) // 2]
        for channel in range(3)
    )
    active = 0
    with warnings.catch_warnings():
        warnings.simplefilter("ignore", DeprecationWarning)
        pixels = list(small.getdata())
    for r, g, b in pixels:
        if abs(r - bg[0]) + abs(g - bg[1]) + abs(b - bg[2]) > 30:
            active += 1
    active_ratio = active / float(small_w * small_h)
    if active_ratio < 0.055:
        return (
            False,
            f"{label} screen region has too little mapped content: active_pixels={active_ratio:.3f}. "
            "This usually means a template screen stayed visible, a crop was pasted incorrectly, "
            "or the captured window is empty.",
        )
    return True, f"{label} ok: mean_luma={mean:.2f}, stddev={stddev:.2f}, active_pixels={active_ratio:.3f}"


def _cross_feature_duplicate_errors(active_names: set[str]) -> list[str]:
    pairs = (
        ("live2d_composite_editor", "mmd_composite_editor"),
        ("live2d_actor_detail", "mmd_character_detail"),
    )
    errors: list[str] = []
    for left, right in pairs:
        if left not in active_names or right not in active_names:
            continue
        left_path = _asset(left)
        right_path = _asset(right)
        if not left_path.exists() or not right_path.exists():
            continue
        try:
            left_hash = _image_sha256(left_path)
            right_hash = _image_sha256(right_path)
        except Exception as exc:
            errors.append(f"- {left} / {right}: could not compare image hashes: {exc}")
            continue
        if left_hash == right_hash:
            errors.append(
                f"- {left} / {right}: {left_path} == {right_path}\n"
                "  Cross-feature duplicate evidence is forbidden. Recapture the target feature "
                "or remove the detail frame instead of reusing another feature's image."
            )
    return errors


def _font(size: int, *, lang: str = "en", mono: bool = False, bold: bool = False) -> ImageFont.ImageFont:
    if mono:
        candidates = [Path("C:/Windows/Fonts/consola.ttf"), Path("C:/Windows/Fonts/cour.ttf")]
    elif lang == "ko":
        candidates = [
            Path("C:/Windows/Fonts/malgun.ttf"),
            Path("C:/Windows/Fonts/malgunbd.ttf"),
            Path("C:/Windows/Fonts/segoeui.ttf"),
        ]
    else:
        candidates = [
            Path("C:/Windows/Fonts/segoeui.ttf"),
            Path("C:/Windows/Fonts/segoeuib.ttf") if bold else Path("C:/Windows/Fonts/segoeui.ttf"),
            Path("C:/Windows/Fonts/arial.ttf"),
        ]
    for path in candidates:
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def _load(path: Path, *, required: bool = True) -> Image.Image:
    if not path.exists():
        raise FileNotFoundError(path)
    return Image.open(path).convert("RGBA")


def _cover(img: Image.Image, size: tuple[int, int]) -> Image.Image:
    target_w, target_h = size
    src_w, src_h = img.size
    scale = max(target_w / src_w, target_h / src_h)
    resized = img.resize((max(1, int(src_w * scale)), max(1, int(src_h * scale))), Image.Resampling.LANCZOS)
    left = max(0, (resized.width - target_w) // 2)
    top = max(0, (resized.height - target_h) // 2)
    return resized.crop((left, top, left + target_w, top + target_h))


def _contain(img: Image.Image, size: tuple[int, int], fill: str = "#101418") -> Image.Image:
    target_w, target_h = size
    src_w, src_h = img.size
    scale = min(target_w / src_w, target_h / src_h)
    resized = img.resize((max(1, int(src_w * scale)), max(1, int(src_h * scale))), Image.Resampling.LANCZOS)
    out = Image.new("RGBA", size, fill)
    out.paste(resized, ((target_w - resized.width) // 2, (target_h - resized.height) // 2), resized if resized.mode == "RGBA" else None)
    return out


def _rounded_paste(
    base: Image.Image,
    img: Image.Image,
    box: tuple[int, int, int, int],
    *,
    radius: int = 6,
    mode: str = "cover",
    fill: str = "#101418",
) -> None:
    x0, y0, x1, y1 = box
    size = (x1 - x0, y1 - y0)
    source = _cover(img.convert("RGBA"), size) if mode == "cover" else _contain(img.convert("RGBA"), size, fill)
    mask = Image.new("L", source.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle((0, 0, source.width, source.height), radius=radius, fill=255)
    base.paste(source, (x0, y0), mask)


def _wrap_text(text: str, width: int, font: ImageFont.ImageFont) -> list[str]:
    lines: list[str] = []
    for raw in text.split("\n"):
        words = raw.split()
        if not words:
            lines.append("")
            continue
        line = ""
        for word in words:
            probe = word if not line else f"{line} {word}"
            if font.getlength(probe) <= width:
                line = probe
            else:
                if line:
                    lines.append(line)
                line = word
        if line:
            lines.append(line)
    return lines


def _clear_left_copy(base: Image.Image) -> None:
    rect = (68, 294, 642, 698)
    fill = Image.new("RGBA", (rect[2] - rect[0], rect[3] - rect[1]), (248, 247, 244, 255))
    mask = Image.new("L", fill.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rectangle((20, 16, fill.width - 20, fill.height - 16), fill=255)
    mask = mask.filter(ImageFilter.GaussianBlur(20))
    base.paste(fill, rect[:2], mask)


def _clear_soft_rect(base: Image.Image, rect: tuple[int, int, int, int], *, blur: int = 8) -> None:
    fill = Image.new("RGBA", (rect[2] - rect[0], rect[3] - rect[1]), (248, 247, 244, 255))
    mask = Image.new("L", fill.size, 0)
    draw = ImageDraw.Draw(mask)
    inset = max(2, blur)
    draw.rectangle((inset, inset, fill.width - inset, fill.height - inset), fill=255)
    mask = mask.filter(ImageFilter.GaussianBlur(blur))
    base.paste(fill, rect[:2], mask)


def _draw_catalog_text(
    base: Image.Image,
    *,
    title: str,
    body: str,
    section: str,
    page: int,
    total: int,
    lang: str,
) -> None:
    draw = ImageDraw.Draw(base)
    _clear_left_copy(base)
    _clear_soft_rect(base, (92, 88, 360, 160), blur=6)
    _clear_soft_rect(base, (88, 806, 520, 862), blur=8)
    _clear_soft_rect(base, (1210, 806, 1585, 862), blur=8)
    mono = _font(14, lang=lang, mono=(lang != "ko"))
    draw.line((115, 105, 272, 105), fill=(142, 142, 138, 255), width=1)
    draw.text((116, 126), section.upper(), fill=(82, 83, 82, 255), font=mono)
    title_font = _font(56 if lang == "en" else 50, lang=lang)
    y = 346
    for line in title.split("\n"):
        draw.text((115, y), line, fill=(31, 32, 33, 255), font=title_font)
        y += 68 if lang == "en" else 62
    body_font = _font(20 if lang == "en" else 19, lang=lang)
    y = max(505, y + 22)
    for line in _wrap_text(body, 420, body_font)[:5]:
        draw.text((118, y), line, fill=(84, 86, 86, 255), font=body_font)
        y += 31 if lang == "en" else 30
    draw.text((115, 607), f"/  {page:02d}  /  {total:02d}", fill=(36, 37, 38, 255), font=_font(16, lang=lang, mono=(lang != "ko")))
    draw.text((115, 831), "TIGERCAPTURE PRODUCT CATALOG", fill=(94, 94, 91, 255), font=mono)
    draw.text((1262, 831), "TIGERCAPTURE.COM", fill=(94, 94, 91, 255), font=mono)
    draw.text((1468, 831), f"PG {page:02d} / {total:02d}", fill=(94, 94, 91, 255), font=mono)


def _label(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str) -> None:
    x, y = xy
    font = _font(17, mono=True)
    pad_x, pad_y = 9, 5
    bbox = draw.textbbox((x, y), text, font=font)
    draw.rounded_rectangle(
        (x - pad_x, y - pad_y, bbox[2] + pad_x, bbox[3] + pad_y),
        radius=7,
        fill=(12, 15, 18, 220),
        outline=(64, 70, 79, 210),
        width=1,
    )
    draw.text((x, y), text, fill=(230, 232, 231, 255), font=font)


def _patch_editor_viewer(editor: Image.Image, frame_path: Path) -> Image.Image:
    patched = editor.copy()
    if frame_path.exists() and editor.width >= 1300 and editor.height >= 800:
        _rounded_paste(patched, _load(frame_path), EDITOR_VIEWER_BOX, radius=10)
    return patched


def _asset(name: str) -> Path:
    paths = {
        "laptop": TEMPLATES / "laptop_catalog_template.png",
        "laptop_ipad": TEMPLATES / "laptop_ipad_catalog_template_v4.png",
        "multi": TEMPLATES / "multi_monitor_front_facing_catalog_template_v2_tight_clean.png",
        "laptop_map": TEMPLATES / "laptop_catalog_template.screen-map.json",
        "laptop_ipad_map": TEMPLATES / "laptop_ipad_catalog_template_v4.screen-map.json",
        "multi_map": TEMPLATES / "multi_monitor_front_facing_catalog_template_v2_tight_clean.screen-map.json",
        "lamborghini_editor": TMP / "fresh_review_recapture" / "catalog_multimedia_lamborghini" / "editor_catalog_multimedia_lamborghini.png",
        "main_editor_current": FRESH_CAPTURE_ROOT / "main_editor_current" / "editor_main_current_timeline_action.png",
        "media_pool_current": FRESH_CAPTURE_ROOT / "media_pool_current" / "editor_media_pool_current_action.png",
        "timeline_current_editor": FRESH_CAPTURE_ROOT / "timeline_current" / "editor_timeline_current_action.png",
        "timeline_current_detail": FRESH_CAPTURE_ROOT / "timeline_current" / "timeline_current_detail_action.png",
        "lamborghini_frame": TMP / "catalog_pretty_frames" / "lamborghini_driving_0126.png",
        "lamborghini_engine": TMP / "catalog_pretty_frames" / "lamborghini_engine_0034.png",
        "taichung": TMP / "catalog_pretty_frames" / "taichung_night_hero_0115.png",
        "tokyo": TMP / "catalog_pretty_frames" / "tokyo_tower_aerial_0223.png",
        "songdo": TMP / "catalog_pretty_frames" / "south_korea_songdo_0924.png",
        "bridge": TMP / "catalog_pretty_frames" / "south_korea_bridge_0351.png",
        "fallingwater": TMP / "catalog_pretty_frames" / "fallingwater_exterior_1328.png",
        "cut_editor": FRESH_CAPTURE_ROOT / "cut_lamborghini" / "editor_cut_edit_action.png",
        "cut_timeline": FRESH_CAPTURE_ROOT / "cut_lamborghini" / "timeline_cut_edit_action.png",
        "overview_left_workspace": FRESH_CAPTURE_ROOT / "multi_environment" / "left_monitor_actor_3d_vtuber_action.png",
        "overview_center_editor": FRESH_CAPTURE_ROOT / "multi_environment" / "center_monitor_editor_action.png",
        "overview_right_workspace": FRESH_CAPTURE_ROOT / "multi_environment" / "right_monitor_node_audio_action.png",
        "effects_before_after_editor": FRESH_CAPTURE_ROOT / "effect_before_after" / "editor_effect_before_after_action.png",
        "effects_hover_detail": FRESH_CAPTURE_ROOT / "effect_before_after" / "effect_hover_or_drag_detail_action.png",
        "effects_editor": FRESH_CAPTURE_ROOT / "effect_southkorea" / "editor_effect_stack_action.png",
        "effects_sheet": FRESH_CAPTURE_ROOT / "effect_southkorea" / "effect_workspace_contact_sheet.png",
        "effects_workbench": FRESH_CAPTURE_ROOT / "effect_southkorea" / "workbench_effect_stack_action.png",
        "transitions_editor": FRESH_CAPTURE_ROOT / "transition_between_clips" / "editor_transition_between_clips_action.png",
        "transition_detail": FRESH_CAPTURE_ROOT / "transition_between_clips" / "transition_timeline_detail_action.png",
        "typography_editor": FRESH_CAPTURE_ROOT / "typography_title_animation" / "editor_typography_title_animation_action.png",
        "typography_detail": FRESH_CAPTURE_ROOT / "typography_title_animation" / "title_animation_detail_action.png",
        "keyframe_motion_editor": FRESH_CAPTURE_ROOT / "keyframe_motion" / "editor_keyframe_motion_action.png",
        "keyframe_detail": FRESH_CAPTURE_ROOT / "keyframe_motion" / "keyframe_timeline_detail_action.png",
        "color_before_after_editor": FRESH_CAPTURE_ROOT / "color_before_after" / "editor_color_before_after_action.png",
        "color_before_after_detail": FRESH_CAPTURE_ROOT / "color_before_after" / "color_controls_detail_action.png",
        "node_before_after_editor": FRESH_CAPTURE_ROOT / "node_effect_before_after" / "editor_node_before_after_action.png",
        "node_graph_actual": FRESH_CAPTURE_ROOT / "node_effect_before_after" / "node_graph_actual_action.png",
        "node_effect_before_after_editor": FRESH_CAPTURE_ROOT / "node_effect_library" / "editor_node_effect_before_after_action.png",
        "node_effect_library_detail": FRESH_CAPTURE_ROOT / "node_effect_library" / "node_effect_library_detail_action.png",
        "node_editor": FRESH_CAPTURE_ROOT / "node_color_tokyo" / "editor_workbench_node_graph_action.png",
        "node_graph": FRESH_CAPTURE_ROOT / "node_color_tokyo" / "workbench_node_graph_action.png",
        "color_editor": FRESH_CAPTURE_ROOT / "node_color_tokyo" / "editor_color_dock_action.png",
        "audio_mixer": FRESH_CAPTURE_ROOT / "node_color_tokyo" / "editor_audio_mixer_action.png",
        "titles_panel": FRESH_CAPTURE_ROOT / "node_color_tokyo" / "editor_title_section_open_action.png",
        "transitions_panel": FRESH_CAPTURE_ROOT / "node_color_tokyo" / "editor_transitions_section_open_action.png",
        "ai_editor": FRESH_CAPTURE_ROOT / "ai_workflow" / "editor_workbench_node_graph_action.png",
        "ppt_maker_editor": FRESH_CAPTURE_ROOT / "ppt_maker_timeline_native" / "ppt_maker_timeline_native_action.png",
        "ppt_maker_detail": FRESH_CAPTURE_ROOT / "ppt_maker_timeline_native" / "ppt_maker_detail_action.png",
        "live2d_composite_editor": FRESH_CAPTURE_ROOT / "live2d_actor_composite" / "editor_live2d_actor_composite_action.png",
        "live2d_actor_detail": FRESH_CAPTURE_ROOT / "live2d_actor_composite" / "live2d_actor_detail_action.png",
        "live2d_editor": FRESH_CAPTURE_ROOT / "live2d_simple_bg" / "editor_live2d_actor_action.png",
        "live2d_viewer": FRESH_CAPTURE_ROOT / "live2d_simple_bg" / "live2d_viewer_action.png",
        "live2d_workbench": FRESH_CAPTURE_ROOT / "live2d_simple_bg" / "workbench_live2d_actor_action.png",
        "mmd_composite_editor": FRESH_CAPTURE_ROOT / "mmd_character_composite" / "editor_mmd_character_composite_action.png",
        "mmd_character_detail": FRESH_CAPTURE_ROOT / "mmd_character_composite" / "mmd_character_detail_action.png",
        "mmd_viewer": FRESH_CAPTURE_ROOT / "mmd_character_motion" / "mmd_player_cantarella_action.png",
        "sound_editor": FRESH_CAPTURE_ROOT / "audio_workbench" / "editor_sound_editor_action.png",
        "sound_workbench": FRESH_CAPTURE_ROOT / "audio_workbench" / "workbench_sound_editor_action.png",
        "sound_graphs": FRESH_CAPTURE_ROOT / "audio_workbench" / "sound_editor_graphs_contact_sheet.png",
        "sound_eq": FRESH_CAPTURE_ROOT / "audio_workbench" / "sound_editor_graph_eq.png",
        "sound_dyn": FRESH_CAPTURE_ROOT / "audio_workbench" / "sound_editor_graph_dyn.png",
        "sound_fx": FRESH_CAPTURE_ROOT / "audio_workbench" / "sound_editor_graph_fx.png",
        "ar_editor": FRESH_CAPTURE_ROOT / "ar_pbr_statue_composite" / "editor_ar_pbr_object_action.png",
        "ar_statue_editor": AR_PBR_MATCHED_EDITOR_CAPTURE,
        "ar_composite": FRESH_CAPTURE_ROOT / "ar_pbr_statue_composite" / "viewer_ar_pbr_composited_frame.png",
        "ar_statue": FRESH_CAPTURE_ROOT / "ar_pbr_statue_composite" / "ar_pbr_statue_standalone_action.png",
        "ar_camera": FRESH_CAPTURE_ROOT / "ar_pbr_camera" / "polyhaven_camera_3d_viewer_no_cubemap_actual.png",
        "ar_car": FRESH_CAPTURE_ROOT / "ar_pbr_nexus_rx" / "nexus_rx_preview_action_zoom.png",
        "vtuber_studio_editor": FRESH_CAPTURE_ROOT / "vrm_vtuber_studio" / "vtuber_broadcast_studio_action.png",
        "vtuber_studio_program_output": FRESH_CAPTURE_ROOT / "vrm_vtuber_studio" / "vtuber_program_output_action.png",
        "vtuber_studio_tracking_mapping": FRESH_CAPTURE_ROOT / "vrm_vtuber_studio" / "vtuber_tracking_mapping_detail_action.png",
        "vtuber_studio_avatar_mapping": FRESH_CAPTURE_ROOT / "vrm_vtuber_studio" / "vtuber_avatar_mapping_detail_action.png",
        "vrm_milica": FRESH_CAPTURE_ROOT / "vrm_vtuber" / "milica_vrm_viewer_action.png",
        "export_editor_current": FRESH_CAPTURE_ROOT / "export_render_queue_current" / "editor_export_render_queue_action.png",
        "export_timeline_detail_current": FRESH_CAPTURE_ROOT / "export_render_queue_current" / "timeline_export_detail_action.png",
        "timeline_visual": FRESH_CAPTURE_ROOT / "export_render_queue" / "timeline_visual_alignment.png",
    }
    return paths[name]


def _vtuber_capture_contract_path() -> Path:
    return FRESH_CAPTURE_ROOT / "vrm_vtuber_studio" / "vtuber_capture_contract.json"


def _vtuber_capture_contract_is_ready(path: Path | None = None) -> tuple[bool, str]:
    contract_path = path or _vtuber_capture_contract_path()
    if not contract_path.exists():
        return (
            False,
            "Missing VTuber capture contract. Recapture with "
            "tools/capture_review_vtuber_studio.py so the build can verify "
            "upper-body VRM evidence instead of trusting image presence only.",
        )
    try:
        data = json.loads(contract_path.read_text(encoding="utf-8"))
    except Exception as exc:
        return False, f"Could not read VTuber capture contract: {exc}"

    evidence = data.get("avatar_evidence") if isinstance(data.get("avatar_evidence"), dict) else {}
    required = set(evidence.get("minimum_visible_parts") or ["head", "neck", "shoulders", "upper_torso"])
    visible = set(evidence.get("visible_parts") or [])
    required_state = data.get("required_state") if isinstance(data.get("required_state"), dict) else {}
    visual_source = str(evidence.get("visual_source") or required_state.get("visual_source") or "")
    renderer_family = str(evidence.get("renderer_family") or required_state.get("renderer_family") or "")
    render_profile = str(evidence.get("render_profile") or required_state.get("render_profile") or "")
    renderer = str(
        evidence.get("renderer_backend")
        or evidence.get("renderer")
        or required_state.get("renderer_backend")
        or required_state.get("renderer")
        or ""
    )

    if visual_source == "vrm_meta_thumbnail_texture":
        return False, "VTuber evidence uses the VRM meta thumbnail; face-only thumbnails are invalid for Trump upper-body source."
    if not bool(evidence.get("review_product_evidence")):
        return False, "VTuber avatar evidence is not marked as product-valid upper-body evidence."
    missing = sorted(required - visible)
    if missing:
        return False, "VTuber avatar evidence is missing visible parts: " + ", ".join(missing)
    if renderer_family != "vtuber_vrm" or render_profile != "vrm_mtoon":
        return False, f"VTuber renderer boundary violation: family={renderer_family!r}, profile={render_profile!r}"
    if renderer != "vrm_mtoon_gpu" or (
        bool(evidence.get("gpu_renderer_required")) and not bool(evidence.get("gpu_renderer_used"))
    ):
        return False, (
            "VTuber catalog evidence must use the VTuber VRM GPU renderer "
            f"(`vrm_mtoon_gpu`), not software/dotted fallback output. renderer={renderer!r}"
        )
    if "software" in renderer.casefold() or "software" in visual_source.casefold():
        return False, "VTuber evidence reports a software renderer/source; product catalog requires GPU VRM output."
    if bool(evidence.get("ar_pbr_used")) or bool(evidence.get("pbr_used")):
        return False, "VTuber evidence must not use AR/PBR, Marmoset PBR, or full-gpu renderer paths."
    return True, "VTuber upper-body evidence contract ready"


def _contract_string_values(value: object) -> list[str]:
    if isinstance(value, str):
        return [value]
    if isinstance(value, dict):
        values: list[str] = []
        for key, item in value.items():
            values.append(str(key))
            values.extend(_contract_string_values(item))
        return values
    if isinstance(value, (list, tuple, set)):
        values = []
        for item in value:
            values.extend(_contract_string_values(item))
        return values
    return []


def _same_file_path(left: Path, right: Path) -> bool:
    try:
        return left.resolve() == right.resolve()
    except Exception:
        return str(left).replace("\\", "/").casefold() == str(right).replace("\\", "/").casefold()


def _vtuber_asset_contract_is_ready(
    name: str,
    image_path: Path,
    *,
    contract_path: Path | None = None,
) -> tuple[bool, str]:
    if name not in VTUBER_VALIDATED_ASSETS:
        return True, "not a VTuber-validated asset"

    resolved_contract = contract_path or _vtuber_capture_contract_path()
    ok, reason = _vtuber_capture_contract_is_ready(resolved_contract)
    if not ok:
        return False, reason
    try:
        data = json.loads(resolved_contract.read_text(encoding="utf-8"))
    except Exception as exc:
        return False, f"Could not read VTuber capture contract: {exc}"

    strings = [str(image_path), *_contract_string_values(data)]
    combined = "\n".join(strings).replace("\\", "/").casefold()
    forbidden = [marker for marker in VTUBER_FORBIDDEN_EVIDENCE_MARKERS if marker.casefold() in combined]
    if forbidden:
        return (
            False,
            "VTuber slide asset is contaminated with forbidden non-VTuber evidence markers "
            f"({', '.join(forbidden)}). Recapture VTuber Studio; do not reuse AR/PBR camera/Marmoset/full-gpu proof.",
        )

    inputs = data.get("inputs") if isinstance(data.get("inputs"), dict) else {}
    for key, marker in VTUBER_REQUIRED_INPUT_MARKERS.items():
        value = str(inputs.get(key) or "").replace("\\", "/").casefold()
        if marker.casefold() not in value:
            return False, f"VTuber capture contract input {key!r} does not match required source marker {marker!r}."

    outputs = data.get("catalog_outputs") if isinstance(data.get("catalog_outputs"), dict) else {}
    raw_output = str(outputs.get(name) or "").strip()
    if not raw_output:
        return False, f"VTuber asset {name!r} is not bound in catalog_outputs. Recapture with catalog_out_dir."
    output_path = Path(raw_output)
    if not output_path.is_absolute():
        output_path = resolved_contract.parent / output_path
    if not _same_file_path(image_path, output_path):
        return False, f"VTuber asset {name!r} points to {output_path}, not the slide source {image_path}."

    hashes = data.get("catalog_output_sha256") if isinstance(data.get("catalog_output_sha256"), dict) else {}
    expected_hash = str(hashes.get(name) or "").strip().casefold()
    if not expected_hash:
        return (
            False,
            f"VTuber asset {name!r} has no catalog_output_sha256. Recapture so camera/AR substitutions cannot reuse the filename.",
        )
    actual_hash = _image_sha256(image_path).casefold()
    if actual_hash != expected_hash:
        return (
            False,
            f"VTuber asset {name!r} hash does not match its capture contract. "
            "The PNG was replaced or came from a stale/non-VTuber source.",
        )

    if name == "vtuber_studio_program_output" and "program_output" not in image_path.name.casefold():
        return False, "VTuber Program Output slide asset must be the Program Output crop, not a tracking/mapping/editor crop."
    if name == "vtuber_studio_avatar_mapping" and "avatar_mapping" not in image_path.name.casefold():
        return False, "VTuber Avatar Mapping slide asset must be the Avatar Mapping crop."

    return True, f"VTuber asset {name} is bound to verified catalog output"


def _page_asset_names() -> set[str]:
    names = {
        "multi",
        "multi_map",
        "laptop",
        "laptop_map",
        "laptop_ipad",
        "laptop_ipad_map",
        "overview_left_workspace",
        "overview_center_editor",
        "overview_right_workspace",
        "ar_statue",
    }
    for spec in PAGES:
        names.add(spec.laptop_name)
        if spec.laptop_frame:
            names.add(spec.laptop_frame)
        if spec.ipad_name:
            names.add(spec.ipad_name)
    return names


def _preflight_required_assets() -> None:
    required_reasons = {
        "ar_statue_editor": AR_PBR_SAME_ASSET_RULE,
        "overview_left_workspace": MULTI_MONITOR_OVERVIEW_RULE,
        "overview_center_editor": MULTI_MONITOR_OVERVIEW_RULE,
        "overview_right_workspace": MULTI_MONITOR_OVERVIEW_RULE,
        "main_editor_current": FEATURE_EVIDENCE_RULE,
        "media_pool_current": FEATURE_EVIDENCE_RULE,
        "timeline_current_editor": FEATURE_EVIDENCE_RULE,
        "timeline_current_detail": (
            "Current timeline evidence lock: recapture from the current editor "
            "timeline visual system. Do not use older V1/A1 block tabs, "
            "synthetic strips, or obsolete thumbnail layouts."
        ),
        "effects_before_after_editor": COMPARE_EVIDENCE_RULE,
        "effects_hover_detail": FEATURE_EVIDENCE_RULE,
        "transitions_editor": FEATURE_EVIDENCE_RULE,
        "transition_detail": FEATURE_EVIDENCE_RULE,
        "typography_editor": FEATURE_EVIDENCE_RULE,
        "typography_detail": FEATURE_EVIDENCE_RULE,
        "keyframe_motion_editor": FEATURE_EVIDENCE_RULE,
        "keyframe_detail": FEATURE_EVIDENCE_RULE,
        "color_before_after_editor": COMPARE_EVIDENCE_RULE,
        "color_before_after_detail": COLOR_IPAD_DETAIL_RULE,
        "ppt_maker_editor": (
            "PPT Maker evidence lock: capture the real app/pptgen or PPT Maker "
            ".tgppt editing surface with video_actor, typography/text, chart/table "
            "or action cards, AR/PBR actor material when available, and timeline "
            "clip bars. Do not substitute a generic PowerPoint mockup or stale "
            "debugCapture screenshot."
        ),
        "ppt_maker_detail": (
            "PPT Maker detail lock: capture a focused PPT Maker detail such as "
            "ppt.* actions, export/snapshot controls, validation/contact sheet, "
            "or selected element inspector. Do not duplicate the whole laptop "
            "screen as the iPad/detail source."
        ),
        "node_before_after_editor": COMPARE_EVIDENCE_RULE,
        "node_graph_actual": COMPARE_EVIDENCE_RULE,
        "node_effect_before_after_editor": COMPARE_EVIDENCE_RULE,
        "node_effect_library_detail": COMPARE_EVIDENCE_RULE,
        "live2d_composite_editor": ACTOR_COMPOSITE_RULE,
        "live2d_actor_detail": ACTOR_COMPOSITE_RULE,
        "vtuber_studio_editor": VTUBER_STUDIO_RULE,
        "vtuber_studio_program_output": VTUBER_STUDIO_RULE,
        "mmd_composite_editor": ACTOR_COMPOSITE_RULE,
        "mmd_character_detail": ACTOR_COMPOSITE_RULE,
        "export_editor_current": FEATURE_EVIDENCE_RULE,
        "export_timeline_detail_current": FEATURE_EVIDENCE_RULE,
    }
    missing: list[str] = []
    forbidden: list[str] = []
    invalid: list[str] = []
    active_names = _page_asset_names()
    for name in sorted(active_names):
        path = _asset(name)
        normalized = str(path).replace("\\", "/").lower()
        if any(marker in normalized for marker in FORBIDDEN_FINAL_CAPTURE_MARKERS):
            forbidden.append(f"- {name}: {path}")
        if not path.exists():
            reason = required_reasons.get(
                name,
                "Required current product-catalog source is missing. Recapture or restage it; do not substitute an older screenshot.",
            )
            missing.append(f"- {name}: {path}\n  {reason}")
            continue
        if name in VIEWER_VALIDATED_ASSETS:
            ok, reason = _editor_viewer_region_is_catalog_ready(path, asset_name=name)
            if not ok:
                invalid.append(f"- {name}: {path}\n  {reason}")
        if name in COMPARE_VALIDATED_ASSETS:
            ok, reason = _compare_capture_contract_is_ready(name, path)
            if not ok:
                invalid.append(f"- {name}: {path}\n  {reason}")
        if name in VTUBER_VALIDATED_ASSETS:
            ok, reason = _vtuber_asset_contract_is_ready(name, path)
            if not ok:
                invalid.append(f"- {name}: {path}\n  {reason}")
        ok, reason = _semantic_capture_contract_is_ready(name, path)
        if not ok:
            invalid.append(f"- {name}: {path}\n  {reason}")
        ok, reason = _semantic_capture_visual_is_ready(name, path)
        if not ok:
            invalid.append(f"- {name}: {path}\n  {reason}")
    ok, reason = _vtuber_capture_contract_is_ready()
    if not ok:
        invalid.append(f"- vtuber_studio.contract: {_vtuber_capture_contract_path()}\n  {reason}")
    invalid.extend(_cross_feature_duplicate_errors(active_names))
    for spec in PAGES:
        if spec.kind == "multi" or spec.key == "closing":
            continue
        if spec.ipad_contract and spec.ipad_name:
            ipad_path = _asset(spec.ipad_name)
            if ipad_path.exists():
                ok, reason = _ipad_detail_contract_is_ready(spec, ipad_path)
                if not ok:
                    invalid.append(f"- {spec.key}.ipad: {ipad_path}\n  {reason}")
        path = _asset(spec.laptop_name)
        if not path.exists():
            continue
        try:
            with Image.open(path) as img:
                if img.width < 900 or img.height < 520:
                    invalid.append(
                        f"- {spec.key}.laptop: {path}\n"
                        f"  Laptop screen source must be a full editor/window capture, not a cropped detail panel: {img.width}x{img.height}"
                    )
        except Exception as exc:
            invalid.append(f"- {spec.key}.laptop: {path}\n  Could not inspect laptop source image: {exc}")
    if forbidden:
        raise RuntimeError(
            "Forbidden historical screenshot source in product-catalog PPT asset map:\n"
            + "\n".join(forbidden)
            + "\nUse current captures under "
            + str(FRESH_CAPTURE_ROOT)
            + " instead. Do not fall back to old screenshots."
        )
    report_lines: list[str] = []
    if missing:
        report_lines.append("Missing required product-catalog capture assets:")
        report_lines.extend(missing)
    if invalid:
        report_lines.append("Invalid product-catalog capture assets:")
        report_lines.extend(invalid)
        report_lines.append("Recapture these pages from the live editor state. Do not patch in raw source frames.")
    if report_lines:
        STRICT_REPORT.parent.mkdir(parents=True, exist_ok=True)
        STRICT_REPORT.write_text("\n".join(report_lines), encoding="utf-8")
        BUILD_BLOCKED_REPORT.parent.mkdir(parents=True, exist_ok=True)
        BUILD_BLOCKED_REPORT.write_text(
            "# Full Product Catalog Build Blocked\n\n"
            "The deck was not generated because one or more current product-catalog "
            "captures are missing or invalid. This prevents repeated placeholder "
            "slides, stale screenshots, and cropped detail panels from being used "
            "as laptop/monitor evidence.\n\n"
            f"Strict report:\n\n```text\n{STRICT_REPORT}\n```\n",
            encoding="utf-8",
        )
        BUILD_WARNINGS.append(str(STRICT_REPORT))
        raise RuntimeError(
            "Full product catalog build blocked by missing/invalid current captures. "
            f"See strict report: {STRICT_REPORT}"
        )
    elif STRICT_REPORT.exists():
        STRICT_REPORT.unlink()


def _editor_viewer_region_is_catalog_ready(path: Path, *, asset_name: str = "") -> tuple[bool, str]:
    try:
        img = Image.open(path).convert("RGB")
    except Exception as exc:
        return False, f"Could not open image for viewer validation: {exc}"
    if img.width < 900 or img.height < 520:
        return False, f"Image is too small for a catalog editor capture: {img.width}x{img.height}"
    ref_w, ref_h = EDITOR_REFERENCE_SIZE
    x0, y0, x1, y1 = EDITOR_VIEWER_BOX_BY_ASSET.get(asset_name, EDITOR_VIEWER_BOX)
    sx = img.width / float(ref_w)
    sy = img.height / float(ref_h)
    box = (
        max(0, min(img.width - 2, int(round(x0 * sx)))),
        max(0, min(img.height - 2, int(round(y0 * sy)))),
        max(1, min(img.width, int(round(x1 * sx)))),
        max(1, min(img.height, int(round(y1 * sy)))),
    )
    if box[2] <= box[0] or box[3] <= box[1]:
        return False, f"Could not derive editor viewer validation box from {img.width}x{img.height}"
    gray = img.crop(box).convert("L")
    stat = ImageStat.Stat(gray)
    mean = float(stat.mean[0])
    stddev = float(stat.stddev[0])
    if mean < 8.0 or stddev < 3.0:
        return (
            False,
            f"Viewer region appears blank/black: mean_luma={mean:.2f}, stddev={stddev:.2f}, box={box}",
        )
    return True, f"viewer ok: mean_luma={mean:.2f}, stddev={stddev:.2f}, box={box}"


def _ipad_detail_contract_is_ready(spec: "PageSpec", path: Path) -> tuple[bool, str]:
    contract = spec.ipad_contract
    if not contract:
        return True, "no iPad detail contract"
    if contract != "color_controls_only":
        return False, f"Unknown iPad detail contract: {contract}"

    name = path.name.lower()
    if any(marker in name for marker in DETAIL_IPAD_FORBIDDEN_NAME_MARKERS):
        return (
            False,
            "Color-grading iPad/detail source must not be a full editor, viewer, "
            f"timeline, or media-pool capture. File name is blocked by contract: {path.name}",
        )
    if not any(marker in name for marker in COLOR_IPAD_ALLOWED_NAME_MARKERS):
        return (
            False,
            "Color-grading iPad/detail source must be a dedicated controls-only "
            "capture named with one of: "
            + ", ".join(COLOR_IPAD_ALLOWED_NAME_MARKERS)
            + f". Current file: {path.name}",
        )
    try:
        with Image.open(path) as img:
            width, height = img.size
    except Exception as exc:
        return False, f"Could not inspect iPad/detail source image: {exc}"

    min_w, min_h = DETAIL_IPAD_MIN_SIZE
    if width < min_w or height < min_h:
        return (
            False,
            f"Color-grading iPad/detail source is too small to explain the controls: {width}x{height}",
        )
    full_w, full_h = DETAIL_IPAD_FULL_EDITOR_MIN_SIZE
    if width >= full_w and height >= full_h:
        return (
            False,
            "Color-grading iPad/detail source looks like a full editor capture "
            f"({width}x{height}). Crop or capture only color wheels, curves, scopes, "
            "and sliders. The iPad must not contain the video viewer, timeline, or media pool.",
        )
    return True, f"{contract} ok: {width}x{height}"


def _editor(name: str, frame: str | None = None) -> Image.Image:
    img = _load(_asset(name))
    if frame:
        return _patch_editor_viewer(img, _asset(frame))
    return img


def _screen_map(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def _screen_rect(data: dict, key: str) -> dict:
    if "screens" in data and key in data["screens"]:
        return data["screens"][key]
    if "regions" in data and key in data["regions"]:
        region = data["regions"][key]
        return region.get("rect", region)
    for region in data.get("screen_regions", ()):
        if region.get("id") == key:
            rect = region.get("rect", region)
            return rect
    if key == "laptop_screen" and "screen" in data:
        return data["screen"]
    raise KeyError(f"screen region not found in template map: {key}")


def _paste_laptop_ipad(base: Image.Image, laptop: Image.Image, ipad: Image.Image) -> None:
    data = _screen_map(_asset("laptop_ipad_map"))
    for key, source in {"laptop_screen": laptop, "ipad_screen": ipad}.items():
        rect = _screen_rect(data, key)
        fitted = _cover(source.convert("RGBA"), (rect["width"], rect["height"]))
        base.paste(fitted, (rect["x"], rect["y"]))


def _paste_laptop(base: Image.Image, laptop: Image.Image) -> None:
    data = _screen_map(_asset("laptop_map"))
    rect = _screen_rect(data, "laptop_screen")
    fitted = _cover(laptop.convert("RGBA"), (rect["width"], rect["height"]))
    base.paste(fitted, (rect["x"], rect["y"]))


def _rect_tuple(rect: dict) -> tuple[int, int, int, int]:
    x = int(rect["x"])
    y = int(rect["y"])
    return x, y, x + int(rect["width"]), y + int(rect["height"])


def _slide_screen_regions(spec: "PageSpec") -> list[tuple[str, tuple[int, int, int, int]]]:
    if spec.kind == "multi":
        data = _screen_map(_asset("multi_map"))
        return [
            (str(region["id"]), _rect_tuple(region["rect"]))
            for region in data.get("screen_regions", ())
        ]
    if spec.key == "closing":
        return []
    if spec.uses_ipad():
        data = _screen_map(_asset("laptop_ipad_map"))
        return [
            ("laptop_screen", _rect_tuple(_screen_rect(data, "laptop_screen"))),
            ("ipad_screen", _rect_tuple(_screen_rect(data, "ipad_screen"))),
        ]
    data = _screen_map(_asset("laptop_map"))
    return [("laptop_screen", _rect_tuple(_screen_rect(data, "laptop_screen")))]


def _crop(img: Image.Image, box: tuple[int, int, int, int], *, contain: bool = False) -> Image.Image:
    cropped = img.crop(box)
    return _contain(cropped, (900, 600)) if contain else cropped


def _make_monitor_left() -> Image.Image:
    return _load(_asset("overview_left_workspace"))


def _make_monitor_center() -> Image.Image:
    return _load(_asset("overview_center_editor"))


def _make_monitor_right() -> Image.Image:
    return _load(_asset("overview_right_workspace"))


def _make_multi_monitor_slide(spec: "PageSpec", lang: str, page: int, total: int, out_path: Path) -> None:
    base = _load(_asset("multi"))
    data = _screen_map(_asset("multi_map"))
    sources = {
        "left_monitor": _make_monitor_left(),
        "center_monitor": _make_monitor_center(),
        "right_monitor": _make_monitor_right(),
    }
    for region in data["screen_regions"]:
        rect = region["rect"]
        fitted = _cover(sources[region["id"]].convert("RGBA"), (rect["width"], rect["height"]))
        base.paste(fitted, (rect["x"], rect["y"]))
    _draw_catalog_text(base, title=spec.title(lang), body=spec.body(lang), section=spec.section(lang), page=page, total=total, lang=lang)
    base.save(out_path)


def _make_feature_slide(spec: "PageSpec", lang: str, page: int, total: int, out_path: Path) -> None:
    laptop = spec.laptop()
    if spec.uses_ipad():
        base = _load(_asset("laptop_ipad"))
        ipad = spec.ipad()
        _paste_laptop_ipad(base, laptop, ipad)
    else:
        base = _load(_asset("laptop"))
        _paste_laptop(base, laptop)
    _draw_catalog_text(base, title=spec.title(lang), body=spec.body(lang), section=spec.section(lang), page=page, total=total, lang=lang)
    base.save(out_path)


def _make_spec_closing_slide(spec: "PageSpec", lang: str, page: int, total: int, out_path: Path) -> None:
    base = Image.new("RGBA", (SLIDE_W, SLIDE_H), "#f8f7f4")
    draw = ImageDraw.Draw(base, "RGBA")
    mono = _font(16, mono=True)
    title_font = _font(56, lang=lang)
    body_font = _font(13, lang=lang)
    small_font = _font(11, lang=lang)
    section = "SPECIFICATION INDEX" if lang == "en" else "스펙 인덱스"
    title = "Specification\nIndex" if lang == "en" else "스펙\n인덱스"
    subtitle = (
        "A compact map of the studio surface: capture, edit, actors, 3D, AI, audio, color, delivery."
        if lang == "en"
        else "캡처, 편집, 액터, 3D, AI, 오디오, 컬러, 전달까지 한 화면에서 이어지는 스튜디오 구성."
    )
    draw.line((115, 108, 272, 108), fill=(108, 108, 104, 255), width=1)
    draw.text((115, 128), section, fill=(86, 87, 84, 255), font=mono)
    draw.text((115, 348), title, fill=(29, 30, 31, 255), font=title_font, spacing=8)
    subtitle_font = _font(18, lang=lang)
    subtitle_y = 518
    for line in _wrap_text(subtitle, 430, subtitle_font)[:3]:
        draw.text((118, subtitle_y), line, fill=(93, 94, 91, 255), font=subtitle_font)
        subtitle_y += 27
    draw.text((115, 831), "TIGERCAPTURE PRODUCT CATALOG", fill=(94, 94, 91, 255), font=mono)
    draw.text((1262, 831), "TIGERCAPTURE.COM", fill=(94, 94, 91, 255), font=mono)
    draw.text((1468, 831), f"PG {page:02d} / {total:02d}", fill=(94, 94, 91, 255), font=mono)
    draw.line((115, 817, 490, 817), fill=(164, 164, 158, 255), width=1)

    if SPEC_CLOSING_BONSAI.exists():
        bonsai = Image.open(SPEC_CLOSING_BONSAI).convert("RGBA")
        bonsai.thumbnail((440, 500), Image.Resampling.LANCZOS)
        bx, by = 1130, 270
        alpha = bonsai.getchannel("A")
        bbox = alpha.getbbox()
        if bbox:
            bottom_band = alpha.crop((0, max(0, bbox[3] - max(12, bonsai.height // 18)), bonsai.width, bbox[3]))
            band_bbox = bottom_band.getbbox()
            if band_bbox:
                sx0, _, sx1, _ = band_bbox
            else:
                sx0, sx1 = bbox[0], bbox[2]
            shadow_w = max(80, int((sx1 - sx0) * 0.82))
            shadow_h = max(10, int(shadow_w * 0.08))
            shadow_x = bx + (sx0 + sx1 - shadow_w) // 2
            shadow_y = by + bbox[3] - max(4, shadow_h // 3)
            shadow = Image.new("RGBA", (shadow_w + 32, shadow_h + 24), (0, 0, 0, 0))
            shadow_draw = ImageDraw.Draw(shadow, "RGBA")
            shadow_draw.ellipse(
                (16, 6, 16 + shadow_w, 6 + shadow_h),
                fill=(46, 42, 36, 52),
            )
            shadow = shadow.filter(ImageFilter.GaussianBlur(7))
            base.alpha_composite(shadow, (shadow_x - 16, shadow_y - 6))
        base.paste(bonsai, (bx, by), bonsai)

    if lang == "ko":
        groups = [
            ("캡처 / 미디어", ["화면 녹화", "YouTube Imports", "미디어 풀", "프록시/썸네일", "다중 포맷"]),
            ("편집 코어", ["타임라인", "컷/분할", "멀티 트랙", "트랜지션", "키프레임"]),
            ("피니싱", ["컬러 그레이딩", "노드 이펙트", "마스크", "비포/애프터", "렌더 큐"]),
            ("캐릭터", ["Live2D", "Spine", "VRM VTuber", "MMD", "액터 트랙"]),
            ("3D / 합성", ["AR/PBR", "Depth-aware", "오클루전", "Shadow Catch", "실시간 조명"]),
            ("AI / 자동화", ["로컬 AI", "Claude", "Python Action", "MCP", "시나리오 실행"]),
            ("오디오", ["사운드 에디터", "EQ", "Dynamics", "FX Curves", "레벨 그래프"]),
            ("제품 기반", ["다국어 UI", "도킹 창", "멀티 모니터", "카탈로그 캡처", "확장 워크플로"]),
        ]
    else:
        groups = [
            ("Capture / Media", ["Screen recording", "YouTube Imports", "Media Pool", "Proxy/thumbnailing", "Multi-format"]),
            ("Editing Core", ["Timeline", "Cut / split", "Multi-track", "Transitions", "Keyframes"]),
            ("Finishing", ["Color grading", "Node effects", "Masks", "Before / after", "Render queue"]),
            ("Actors", ["Live2D", "Spine", "VRM VTuber", "MMD", "Actor lanes"]),
            ("3D / Composite", ["AR/PBR", "Depth-aware", "Occlusion", "Shadow catch", "Realtime lighting"]),
            ("AI / Automation", ["Local AI", "Claude", "Python Action", "MCP", "Scenario runs"]),
            ("Audio", ["Sound editor", "EQ", "Dynamics", "FX curves", "Level graphs"]),
            ("Product Surface", ["Multilingual UI", "Dockable windows", "Multi-monitor", "Catalog capture", "Extensible workflows"]),
        ]
    col_x = [610, 790, 970]
    y0 = 170
    row_h = 160
    for idx, (heading, items) in enumerate(groups):
        x = col_x[idx % 3]
        y = y0 + (idx // 3) * row_h
        draw.text((x, y), heading, fill=(30, 31, 32, 255), font=body_font)
        yy = y + 25
        for item in items:
            draw.text((x, yy), f"- {item}", fill=(82, 84, 82, 255), font=small_font)
            yy += 19
    base.save(out_path)


@dataclass(frozen=True)
class PageSpec:
    key: str
    section_en: str
    title_en: str
    body_en: str
    section_ko: str
    title_ko: str
    body_ko: str
    laptop_name: str = "lamborghini_editor"
    laptop_frame: str | None = "lamborghini_frame"
    ipad_name: str | None = None
    ipad_crop: tuple[int, int, int, int] | None = None
    ipad_contain: bool = False
    ipad_contract: str | None = None
    kind: str = "feature"

    def title(self, lang: str) -> str:
        return self.title_ko if lang == "ko" else self.title_en

    def body(self, lang: str) -> str:
        return self.body_ko if lang == "ko" else self.body_en

    def section(self, lang: str) -> str:
        return self.section_ko if lang == "ko" else self.section_en

    def laptop(self) -> Image.Image:
        return _editor(self.laptop_name, self.laptop_frame)

    def uses_ipad(self) -> bool:
        return bool(self.ipad_name)

    def ipad(self) -> Image.Image:
        if not self.ipad_name:
            raise RuntimeError(
                f"{self.key} has no iPad/detail source. Use the laptop-only template "
                "instead of duplicating the laptop screen."
            )
        source = _load(_asset(self.ipad_name))
        if self.ipad_crop:
            source = _crop(source, self.ipad_crop, contain=self.ipad_contain)
        elif self.ipad_contain:
            source = _contain(source, (900, 600))
        return source


PAGES = [
    PageSpec(
        "studio_overview",
        "STUDIO OVERVIEW",
        "Multi-Environment\nEditing Studio",
        "Spread the same project across the screens you have: the center stays focused on preview and timeline, while side displays hold actors, 3D, nodes, and audio.",
        "스튜디오 오버뷰",
        "멀티 환경 편집\n스튜디오",
        "보유한 화면 수에 맞춰 같은 프로젝트를 펼칩니다. 센터는 영상 프리뷰와 타임라인에 집중하고, 주변 화면은 액터, 3D, 노드, 사운드를 맡습니다.",
        kind="multi",
    ),
    PageSpec(
        "studio_surface",
        "MAIN EDITOR",
        "TigerCapture\nStudio",
        "Import real footage, arrange tracks, preview the result, and keep media, workbench, AI, and timeline controls in one editing surface.",
        "메인 에디터",
        "TigerCapture\nStudio",
        "실제 영상을 불러오고 트랙에 배치하며, 미디어 풀과 워크벤치, AI, 타임라인을 한 화면에서 다룹니다.",
        laptop_name="main_editor_current",
        laptop_frame=None,
    ),
    PageSpec(
        "ai_workflow",
        "AI WORKFLOW",
        "AI-Driven\nEditing",
        "Claude, Codex, and local LLMs can drive registered editor actions: cuts, filters, speed, tracks, nodes, captures, and scripted workflows.",
        "AI 워크플로",
        "AI 주도\n편집",
        "Claude, Codex, 로컬 LLM이 컷, 필터, 배속, 트랙, 노드, 캡처 같은 등록 액션을 실행합니다.",
        laptop_name="ai_editor",
        laptop_frame="taichung",
        ipad_crop=(360, 520, 1480, 845),
    ),
    PageSpec(
        "ppt_maker",
        "PRESENTATION STUDIO",
        "PPT Maker",
        "Build timeline-native presentations from real media, typography, charts, and 3D actors. Save .tgppt projects, validate the deck, and export PPTX, PNG snapshots, or previews.",
        "프레젠테이션 스튜디오",
        "PPT Maker",
        "실제 미디어, 타이포그래피, 차트, 3D 액터를 하나의 .tgppt 프로젝트에 배치하고 PPTX, PNG 스냅샷, 컨택트 시트, 비디오 프리뷰로 내보냅니다.",
        laptop_name="ppt_maker_editor",
        laptop_frame=None,
    ),
    PageSpec(
        "media_pool",
        "MEDIA",
        "Media Pool\nAnd Imports",
        "YouTube imports, local video, audio, images, Live2D, and 3D assets stay visible as reusable project material.",
        "미디어",
        "미디어 풀과\n불러오기",
        "유튜브 임포트, 로컬 영상, 오디오, 이미지, Live2D, 3D 에셋을 프로젝트 재료로 관리합니다.",
        laptop_name="media_pool_current",
        laptop_frame=None,
        ipad_name="media_pool_current",
        ipad_crop=(0, 0, 250, 520),
    ),
    PageSpec(
        "timeline",
        "EDITING CORE",
        "Timeline\nEditing",
        "Long clips, split points, transitions, markers, and layered tracks keep the edit readable without flattening the production view.",
        "편집 코어",
        "타임라인\n편집",
        "긴 클립, 컷 지점, 트랜지션, 마커, 레이어 트랙을 작업 흐름 안에서 읽기 쉽게 보여줍니다.",
        laptop_name="timeline_current_editor",
        laptop_frame=None,
    ),
    PageSpec(
        "effects",
        "EFFECTS",
        "Drag-First\nEffects",
        "Effects are selected visually, previewed quickly, and dropped into the working edit instead of being hidden behind modal lists.",
        "이펙트",
        "드래그 중심\n이펙트",
        "이펙트를 시각적으로 고르고 미리 확인한 뒤, 모달 창 없이 작업 중인 편집에 바로 배치합니다.",
        laptop_name="effects_before_after_editor",
        laptop_frame=None,
        ipad_name="effects_hover_detail",
        ipad_contain=True,
    ),
    PageSpec(
        "transitions",
        "TRANSITIONS",
        "Transitions\nOn Tracks",
        "Transition presets live beside the timeline so crossfades, wipes, and stylized cuts can be auditioned in context.",
        "트랜지션",
        "트랙 위\n전환 효과",
        "크로스페이드, 와이프, 스타일 컷을 타임라인 옆에서 고르고 실제 장면 맥락으로 확인합니다.",
        laptop_name="transitions_editor",
        laptop_frame=None,
        ipad_name="transition_detail",
        ipad_contain=True,
    ),
    PageSpec(
        "typography",
        "TYPOGRAPHY",
        "Typography\nAnd Titles",
        "Title presets support large text, motion, and timeline placement so captions become designed elements, not afterthoughts.",
        "타이포그래피",
        "타이포그래피와\n타이틀",
        "타이틀 프리셋은 큰 글자, 모션, 타임라인 배치를 지원해 자막을 디자인 요소로 만듭니다.",
        laptop_name="typography_editor",
        laptop_frame=None,
        ipad_name="typography_detail",
        ipad_contain=True,
    ),
    PageSpec(
        "keyframes",
        "MOTION",
        "Keyframes\nAnd Motion",
        "Opacity, scale, movement, and effect values can be animated directly on the timeline for visible editorial timing.",
        "모션",
        "키프레임과\n모션",
        "투명도, 스케일, 이동, 이펙트 값을 타임라인 위에서 애니메이션해 편집 타이밍을 직접 보여줍니다.",
        laptop_name="keyframe_motion_editor",
        laptop_frame=None,
    ),
    PageSpec(
        "color",
        "FINISHING",
        "Color Grading\nWorkspace",
        "Curves, wheels, tone controls, and scopes make color work feel like a finishing room inside the editor.",
        "피니싱",
        "컬러 그레이딩\n워크스페이스",
        "커브, 휠, 톤 컨트롤, 스코프를 통해 편집기 안에서 피니싱룸처럼 색을 다룹니다.",
        laptop_name="color_before_after_editor",
        laptop_frame=None,
        ipad_name="color_before_after_detail",
        ipad_crop=None,
        ipad_contain=True,
        ipad_contract="color_controls_only",
    ),
    PageSpec(
        "node_graph",
        "COMPOSITING",
        "Node Graph\nComposition",
        "Blur, glow, LUT, color grade, mask, blend, overlay, and output nodes can be connected while the timeline remains active.",
        "컴포지팅",
        "노드 그래프\n합성",
        "블러, 글로우, LUT, 컬러 그레이드, 마스크, 블렌드, 오버레이, 출력 노드를 타임라인과 함께 연결합니다.",
        laptop_name="node_before_after_editor",
        laptop_frame=None,
        ipad_name="node_graph_actual",
        ipad_contain=True,
    ),
    PageSpec(
        "node_effects",
        "NODE EFFECTS",
        "Node Effects\nLibrary",
        "Common effect nodes become reusable building blocks, with before/after checks available from the same workbench.",
        "노드 이펙트",
        "노드 이펙트\n라이브러리",
        "자주 쓰는 이펙트 노드를 재사용 가능한 블록으로 두고, 같은 워크벤치에서 비포/애프터를 확인합니다.",
        laptop_name="node_effect_before_after_editor",
        laptop_frame=None,
        ipad_name="node_effect_library_detail",
        ipad_contain=True,
    ),
    PageSpec(
        "audio_workbench",
        "AUDIO",
        "Audio\nWorkbench",
        "The sound editor can sit inside the video workflow, editing selected media or the active audio track without losing prior edits.",
        "오디오",
        "사운드\n워크벤치",
        "사운드 에디터를 영상 워크플로 안에 두고, 선택 미디어나 오디오 트랙을 이전 편집 상태와 함께 다룹니다.",
        laptop_name="sound_editor",
        laptop_frame=None,
        ipad_name="sound_workbench",
        ipad_contain=True,
    ),
    PageSpec(
        "audio_curves",
        "AUDIO DETAIL",
        "EQ, Dynamics\nAnd FX Curves",
        "EQ curves, dynamics, effect automation, and level graphs keep sound design visible instead of burying it in export settings.",
        "오디오 디테일",
        "EQ, 다이내믹스와\nFX 커브",
        "EQ 커브, 다이내믹스, 이펙트 자동화, 레벨 그래프를 작업 화면 안에서 시각적으로 편집합니다.",
        laptop_name="sound_editor",
        laptop_frame=None,
        ipad_name="sound_graphs",
        ipad_contain=True,
    ),
    PageSpec(
        "live2d_spine",
        "ACTOR OVERLAY",
        "Live2D And\nSpine Tracks",
        "Live2D actors sit on dedicated lanes with transform and opacity animation; Spine support is tracked as an actor workflow too.",
        "액터 오버레이",
        "Live2D와\nSpine 트랙",
        "Live2D 액터는 전용 레인에서 이동과 투명도를 애니메이션하고, Spine도 액터 워크플로로 함께 관리합니다.",
        laptop_name="live2d_composite_editor",
        laptop_frame=None,
        ipad_name="live2d_viewer",
        ipad_contain=True,
    ),
    PageSpec(
        "vrm",
        "VTUBER",
        "VRM VTuber\nStudio",
        "VRM sources can become character overlays for creator video, with motion input and presentation-ready character captures.",
        "버튜버",
        "VRM VTuber\nStudio",
        "VRM 소스를 크리에이터 영상용 캐릭터 오버레이로 사용하고, 모션 입력과 발표용 캡처로 연결합니다.",
        laptop_name="vtuber_studio_editor",
        laptop_frame=None,
        ipad_name="vtuber_studio_program_output",
        ipad_contain=True,
    ),
    PageSpec(
        "mmd",
        "CHARACTER MOTION",
        "MMD Character\nMotion",
        "MMD playback and character motion windows can be used beside the main edit for music-video and avatar-driven production.",
        "캐릭터 모션",
        "MMD 캐릭터\n모션",
        "MMD 재생과 캐릭터 모션 창을 메인 편집 옆에 두고 뮤직비디오와 아바타 영상을 제작합니다.",
        laptop_name="mmd_viewer",
        laptop_frame=None,
    ),
    PageSpec(
        "ar_pbr",
        "3D COMPOSITE",
        "AR/PBR\n3D Composite",
        "Real-time PBR assets can be placed over video with depth-aware composition, occlusion, shadow catching, tone, and lighting controls.",
        "3D 합성",
        "AR/PBR\n3D 합성",
        "실시간 PBR 에셋을 영상 위에 놓고 깊이, 오클루전, 그림자, 톤, 조명 컨트롤로 합성합니다.",
        laptop_name="ar_statue_editor",
        laptop_frame=None,
        ipad_name="ar_statue",
        ipad_contain=True,
    ),
    PageSpec(
        "creator_assist",
        "CREATOR ASSIST",
        "Creator Assist\nWorkflow",
        "AI prompts, presets, media choices, and action sequences come together as a guided editing flow for repeatable creator work.",
        "크리에이터 지원",
        "Creator Assist\n워크플로",
        "AI 프롬프트, 프리셋, 미디어 선택, 액션 시퀀스를 반복 가능한 크리에이터 편집 흐름으로 묶습니다.",
        laptop_name="ai_editor",
        laptop_frame="bridge",
        ipad_crop=(445, 525, 1480, 770),
    ),
    PageSpec(
        "export",
        "DELIVERY",
        "Export And\nRender Queue",
        "Delivery keeps preview parity, metadata, tracks, masks, actors, and color decisions aligned through the final render queue.",
        "전달",
        "내보내기와\n렌더 큐",
        "프리뷰 패리티, 메타데이터, 트랙, 마스크, 액터, 컬러 결정을 최종 렌더 큐까지 맞춰 갑니다.",
        laptop_name="export_editor_current",
        laptop_frame=None,
    ),
    PageSpec(
        "closing",
        "SPECIFICATION INDEX",
        "Specification\nIndex",
        "A compact map of TigerCapture's capture, editing, actor, 3D, AI, audio, color, and delivery surface.",
        "스펙 인덱스",
        "스펙\n인덱스",
        "TigerCapture의 캡처, 편집, 액터, 3D, AI, 오디오, 컬러, 전달 기능을 한 페이지에 요약합니다.",
        laptop_name="node_editor",
        laptop_frame="tokyo",
        ipad_name="sound_graphs",
        ipad_contain=True,
    ),
]


def _slug(text: str) -> str:
    return "".join(ch if ch.isalnum() else "_" for ch in text.lower()).strip("_")


def build_slides(lang: str) -> list[Path]:
    slides_dir = OUT / f"slides_{lang}"
    if slides_dir.exists():
        shutil.rmtree(slides_dir)
    slides_dir.mkdir(parents=True, exist_ok=True)
    outputs: list[Path] = []
    total = len(PAGES)
    for index, spec in enumerate(PAGES, start=1):
        out_path = slides_dir / f"{index:02d}_{_slug(spec.key)}.png"
        if spec.key == "closing":
            _make_spec_closing_slide(spec, lang, index, total, out_path)
        elif spec.kind == "multi":
            _make_multi_monitor_slide(spec, lang, index, total, out_path)
        else:
            _make_feature_slide(spec, lang, index, total, out_path)
        outputs.append(out_path)
    return outputs


def build_deck(lang: str, slide_paths: list[Path]) -> Path:
    out_path = _available_output_path(OUT / f"TigerCapture_Product_Catalog_Full_{lang.upper()}.pptx")
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:  # pragma: no cover - dependency guard for local tooling
        raise RuntimeError("python-pptx is required for PowerPoint-compatible catalog export") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for path in slide_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out_path)
    return out_path


def build_contact_sheet(paths: list[Path], out_path: Path, *, cols: int = 3) -> Path:
    thumbs = []
    for path in paths:
        img = _load(path)
        img.thumbnail((480, 270), Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", (480, 270), "#f8f7f4")
        canvas.paste(img.convert("RGB"), ((480 - img.width) // 2, (270 - img.height) // 2))
        thumbs.append(canvas)
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 480, rows * 270), "#222222")
    for idx, thumb in enumerate(thumbs):
        x = (idx % cols) * 480
        y = (idx // cols) * 270
        sheet.paste(thumb, (x, y))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out_path)
    return out_path


def validate_rendered_slides(lang: str, slide_paths: list[Path]) -> None:
    errors: list[str] = []
    report_rows: list[dict[str, object]] = []
    evidence_signatures: dict[str, str] = {}
    if len(slide_paths) != len(PAGES):
        errors.append(f"Rendered slide count mismatch for {lang}: {len(slide_paths)} != {len(PAGES)}")
    for index, (spec, path) in enumerate(zip(PAGES, slide_paths), start=1):
        try:
            img = Image.open(path).convert("RGB")
        except Exception as exc:
            errors.append(f"- slide {index:02d} {spec.key}: could not open rendered PNG: {exc}")
            continue
        if img.size != (SLIDE_W, SLIDE_H):
            errors.append(f"- slide {index:02d} {spec.key}: rendered PNG size mismatch: {img.size} != {(SLIDE_W, SLIDE_H)}")
        gray = img.convert("L")
        stat = ImageStat.Stat(gray)
        mean = float(stat.mean[0])
        stddev = float(stat.stddev[0])
        if mean < 18.0 and stddev < 8.0:
            errors.append(f"- slide {index:02d} {spec.key}: whole slide is too dark/blank: mean={mean:.2f}, stddev={stddev:.2f}")
        if mean > 247.0 and stddev < 9.0:
            errors.append(f"- slide {index:02d} {spec.key}: whole slide is nearly blank white: mean={mean:.2f}, stddev={stddev:.2f}")

        screen_rows: list[dict[str, object]] = []
        crops: dict[str, Image.Image] = {}
        for label, box in _slide_screen_regions(spec):
            crop = img.crop(box)
            crops[label] = crop
            ok, reason = _screen_region_is_catalog_ready(crop, label=f"slide {index:02d} {spec.key}.{label}")
            screen_rows.append({"label": label, "box": box, "ok": ok, "reason": reason})
            if not ok:
                errors.append(f"- slide {index:02d} {spec.key}.{label}: {reason}")
            if label in {"laptop_screen", "center_monitor", "left_monitor", "right_monitor"}:
                signature = _screen_crop_signature(crop)
                previous = evidence_signatures.get(signature)
                if previous and spec.key != "closing":
                    errors.append(
                        f"- slide {index:02d} {spec.key}.{label}: final evidence screen is an exact duplicate of {previous}. "
                        "Recapture the page-specific feature or remove the repeated device frame."
                    )
                else:
                    evidence_signatures[signature] = f"slide {index:02d} {spec.key}.{label}"

        if spec.uses_ipad() and "laptop_screen" in crops and "ipad_screen" in crops:
            laptop_hash = _image_dhash(crops["laptop_screen"])
            ipad_hash = _image_dhash(crops["ipad_screen"])
            distance = _hamming_distance(laptop_hash, ipad_hash)
            if distance <= 3 and spec.ipad_crop is None:
                errors.append(
                    f"- slide {index:02d} {spec.key}: iPad/detail screen is visually the same as the laptop screen "
                    f"(dHash distance {distance}). Use a feature-specific detail crop or the laptop-only template."
                )
        report_rows.append(
            {
                "slide": index,
                "key": spec.key,
                "path": str(path),
                "mean_luma": mean,
                "stddev_luma": stddev,
                "screens": screen_rows,
            }
        )

    report_path = OUT / f"visual_quality_report_{lang}.json"
    report_path.write_text(
        json.dumps(
            {
                "schema": "tigercapture.product_catalog.visual_quality_report.v1",
                "lang": lang,
                "ok": not errors,
                "error_count": len(errors),
                "errors": errors,
                "slides": report_rows,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    if errors:
        BUILD_BLOCKED_REPORT.parent.mkdir(parents=True, exist_ok=True)
        BUILD_BLOCKED_REPORT.write_text(
            "# Full Product Catalog Visual QA Blocked\n\n"
            "The deck was not exported because rendered slide QA found device "
            "mapping, duplicate evidence, or visual emptiness problems.\n\n"
            f"Visual QA report:\n\n```text\n{report_path}\n```\n\n"
            + "\n".join(errors),
            encoding="utf-8",
        )
        raise RuntimeError(f"Rendered slide visual QA failed for {lang}. See {report_path}")


def validate_pptx(path: Path, expected_slides: int) -> None:
    required_entries = {
        "[Content_Types].xml",
        "_rels/.rels",
        "ppt/presentation.xml",
        "ppt/_rels/presentation.xml.rels",
    }
    try:
        with zipfile.ZipFile(path) as zf:
            corrupt = zf.testzip()
            if corrupt:
                raise RuntimeError(f"{path} has a corrupt PPTX zip entry: {corrupt}")
            names = set(zf.namelist())
            missing = sorted(required_entries - names)
            if missing:
                raise RuntimeError(f"{path} is missing required PPTX entries: {', '.join(missing)}")
            xml_entries = [
                name
                for name in zf.namelist()
                if name.endswith(".xml") and (name.startswith("ppt/") or name == "[Content_Types].xml")
            ]
            for name in xml_entries:
                try:
                    ET.fromstring(zf.read(name))
                except Exception as exc:
                    raise RuntimeError(f"{path} contains invalid XML at {name}: {exc}") from exc
            slide_entries = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
            media_entries = [name for name in zf.namelist() if name.startswith("ppt/media/")]
    except zipfile.BadZipFile as exc:
        raise RuntimeError(f"{path} is not a valid PPTX/zip file: {exc}") from exc
    if len(slide_entries) != expected_slides:
        raise RuntimeError(f"{path} slide count mismatch: {len(slide_entries)} != {expected_slides}")
    if len(media_entries) < expected_slides:
        raise RuntimeError(
            f"{path} has fewer embedded media files than slides: {len(media_entries)} < {expected_slides}"
        )
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - dependency guard for local tooling
        raise RuntimeError("python-pptx is required to reopen and validate the exported PPTX") from exc
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise RuntimeError(f"{path} could not be reopened by python-pptx: {exc}") from exc
    if len(prs.slides) != expected_slides:
        raise RuntimeError(f"{path} python-pptx slide count mismatch: {len(prs.slides)} != {expected_slides}")


def _clear_previous_outputs() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    for path in OUT.glob("*.pptx"):
        _safe_unlink(path)
    for path in OUT.glob("contact_sheet_*.png"):
        _safe_unlink(path)
    for path in OUT.glob("candidate_capture_sheet.png"):
        _safe_unlink(path)
    for path in OUT.glob("slides_*"):
        if path.is_dir():
            shutil.rmtree(path)
    if BUILD_BLOCKED_REPORT.exists():
        _safe_unlink(BUILD_BLOCKED_REPORT)


def _safe_unlink(path: Path) -> bool:
    try:
        path.unlink()
        return True
    except FileNotFoundError:
        return True
    except PermissionError:
        BUILD_WARNINGS.append(f"Could not delete locked output; leaving it in place: {path}")
        return False


def _available_output_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem = path.stem
    suffix = path.suffix
    for index in range(2, 100):
        candidate = path.with_name(f"{stem}_{index}{suffix}")
        if not candidate.exists():
            return candidate
    raise RuntimeError(f"Could not find an available output path near {path}")


def main(argv: list[str] | None = None) -> int:
    argv = list(sys.argv[1:] if argv is None else argv)
    if "--preflight-only" in argv:
        _preflight_required_assets()
        print("Full product catalog preflight ok.")
        return 0
    OUT.mkdir(parents=True, exist_ok=True)
    _clear_previous_outputs()
    _preflight_required_assets()
    results = []
    for lang in ("en", "ko"):
        slide_paths = build_slides(lang)
        validate_rendered_slides(lang, slide_paths)
        deck_path = build_deck(lang, slide_paths)
        validate_pptx(deck_path, len(slide_paths))
        sheet = build_contact_sheet(slide_paths, OUT / f"contact_sheet_{lang}.png")
        results.append((lang, deck_path, sheet, slide_paths[0]))
    for lang, deck_path, sheet, first_slide in results:
        print(f"{lang}: {deck_path}")
        print(f"{lang}_contact_sheet: {sheet}")
        print(f"{lang}_first_slide: {first_slide}")
    for warning in BUILD_WARNINGS:
        print(f"strict_report: {warning}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
